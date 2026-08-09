"""
opensecstack EU Conference — dual-language deck generator.
Produces:
  opensecstack_EU_Conference_EN.pptx  (English)
  opensecstack_EU_Conference_AL.pptx  (Albanian / Shqip)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x1A, 0x2F)
TEAL   = RGBColor(0x00, 0xBF, 0xB2)
AMBER  = RGBColor(0xFF, 0xB3, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xE6, 0xEA, 0xF0)
DGREY  = RGBColor(0x44, 0x55, 0x66)
GREEN  = RGBColor(0x00, 0xC8, 0x7A)
RED    = RGBColor(0xFF, 0x33, 0x44)
BLUE   = RGBColor(0x00, 0x66, 0xCC)
PURPLE = RGBColor(0x55, 0x00, 0xCC)
GOLD   = RGBColor(0xFF, 0xD7, 0x00)

W = Inches(13.33)
H = Inches(7.5)

# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.line.width = 0
    return s

def txt(slide, text, x, y, w, h, size=18, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    if not text:
        return
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

# ── Slide templates ───────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle="", footer=""):
    """Full-bleed navy title slide."""
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, H - Inches(0.10), W, Inches(0.10), TEAL)
    rect(s, 0, Inches(0), Inches(0.08), H, TEAL)
    txt(s, title,    Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.0),
        size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, subtitle, Inches(1.0), Inches(4.1), Inches(11.3), Inches(1.0),
        size=22, color=TEAL,  align=PP_ALIGN.CENTER)
    if footer:
        txt(s, footer, Inches(0), H - Inches(0.55), W, Inches(0.45),
            size=13, italic=True, color=LGREY, align=PP_ALIGN.CENTER)
    return s

def slide_section(prs, number, title, subtitle=""):
    """Dark section-break slide."""
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, Inches(3.1), W, Inches(0.06), TEAL)
    txt(s, f"SECTION {number}", Inches(0), Inches(2.3), W, Inches(0.7),
        size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, title,    Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.5),
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, subtitle, Inches(1.0), Inches(4.9), Inches(11.3), Inches(0.8),
        size=20, color=LGREY, align=PP_ALIGN.CENTER, italic=True)
    return s

def slide_content(prs, title):
    """White content slide with navy header."""
    s = blank(prs)
    rect(s, 0, 0, W, H, WHITE)
    rect(s, 0, 0, W, Inches(1.1), NAVY)
    rect(s, 0, Inches(1.1), Inches(0.07), H - Inches(1.1), TEAL)
    txt(s, title, Inches(0.35), Inches(0.15), Inches(12.6), Inches(0.85),
        size=28, bold=True, color=WHITE)
    return s

def bullets(prs, title, items):
    """Bullet list slide."""
    s = slide_content(prs, title)
    y = Inches(1.28)
    for item in items:
        ind  = item.startswith("  ")
        text = item.lstrip()
        pfx  = "    ◦  " if ind else "  ▶  "
        sz   = 17 if ind else 20
        clr  = DGREY if ind else NAVY
        bld  = not ind
        txt(s, pfx + text, Inches(0.4), y, Inches(12.4), Inches(0.5),
            size=sz, bold=bld, color=clr)
        y += Inches(0.50 if not ind else 0.43)
    return s

def box(slide, label, value, x, y, w=Inches(2.8), h=Inches(1.55),
        bg=NAVY, lc=TEAL, vc=WHITE, ls=13, vs=18):
    rect(slide, x, y, w, h, bg)
    txt(slide, label, x + Inches(0.12), y + Inches(0.10),
        w - Inches(0.2), Inches(0.42), size=ls, bold=True, color=lc)
    txt(slide, value, x + Inches(0.12), y + Inches(0.52),
        w - Inches(0.2), h - Inches(0.6), size=vs, color=vc, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# CONTENT DEFINITIONS  —  English + Albanian
# ══════════════════════════════════════════════════════════════════════════════

DECKS = {
    "EN": {
        "filename": "opensecstack_EU_Conference_EN.pptx",
        "title": "opensecstack",
        "subtitle": "Governing the Digital Future — Secure OS · Blockchain Governance · EU Ecosystem",
        "footer": "EU Conference  ·  2026  ·  opensecstack.org",
        "slides": []   # populated below
    },
    "AL": {
        "filename": "opensecstack_EU_Conference_AL.pptx",
        "title": "opensecstack",
        "subtitle": "Qeverisja e së Ardhmes Dixhitale — OS i Sigurt · Qeverisje Blockchain · Ekosistemi BE",
        "footer": "Konferencë BE  ·  2026  ·  opensecstack.org",
        "slides": []
    },
}

# Each entry: ("type", title, [content...])
# types: section | bullets | custom_<name>

EN = DECKS["EN"]["slides"]
AL = DECKS["AL"]["slides"]

# ── Slide 2: Vision ──────────────────────────────────────────────────────────
EN.append(("bullets", "Our Vision: Digital Sovereignty by Design", [
    "Software that cannot be forged, tampered with, or silently bypassed",
    "An open-source ecosystem where every action leaves a cryptographic footprint",
    "An operating system built for security — not bolted on after the fact",
    "Governance that is mathematically provable — not just policy",
    "  Every decision: chain-hashed, anchored, and auditable forever",
    "  Every actor: separated by duty, verified by cryptography",
    "EU-native: NIS2 · GDPR · Cyber Resilience Act — covered by design",
    "Open to all: Apache 2.0 · Multi-language SDK · Community-driven",
]))
AL.append(("bullets", "Vizioni ynë: Sovranitet Dixhital nga Dizajni", [
    "Softuer që nuk mund të falsifikohet, ndryshohet apo anashkalohet",
    "Ekosistem me burim të hapur ku çdo veprim lë gjurmë kriptografike",
    "Sistem operativ i ndërtuar për siguri — jo i shtuar pas faktit",
    "Qeverisje e provueshme matematikisht — jo vetëm politikë",
    "  Çdo vendim: i lidhur me hash zinxhiri, i ankoruar dhe i auditushëm",
    "  Çdo aktor: i ndarë sipas detyrës, i verifikuar me kriptografi",
    "BE-nativ: NIS2 · GDPR · Akti i Rezistencës Kibernetike — i mbuluar nga dizajni",
    "I hapur për të gjithë: Apache 2.0 · SDK shumë-gjuhësh · I drejtuar nga komuniteti",
]))

# ── Slide 3: The Problem ─────────────────────────────────────────────────────
EN.append(("bullets", "The Problem We Solve", [
    "AI systems execute thousands of privileged actions daily — with no verifiable record",
    "  Who authorised this deployment? When? With what evidence?",
    "Audit logs live in databases that administrators can edit",
    "  Post-breach forensics are unreliable — the evidence can be destroyed",
    "Separation of Duties is policy, not enforcement — bypassed silently",
    "  One operator triggers and approves the same change — undetected",
    "Regulatory mandates demand proof, not process",
    "  NIS2 Article 21(2)(h): cryptographic evidence of control is required",
    "Current solutions: expensive, proprietary, single-vendor lock-in",
]))
AL.append(("bullets", "Problemi që Zgjidhim", [
    "Sistemet AI ekzekutojnë mijëra veprime të privilegjuara çdo ditë — pa regjistër të verifikueshëm",
    "  Kush autorizoi këtë vendosje? Kur? Me çfarë prove?",
    "Regjistrat e auditimit ndodhen në baza të dhënash që administratorët mund t'i ndryshojnë",
    "  Hetimi pas shkeljes është i pabesueshëm — provat mund të shkatërrohen",
    "Ndarja e Detyrave është politikë, jo zbatim — anashkalohet në heshtje",
    "  Një operator aktivizon dhe miraton të njëjtin ndryshim — pa u zbuluar",
    "Mandatet rregullatore kërkojnë prova, jo procese",
    "  NIS2 neni 21(2)(h): kërkohen prova kriptografike të kontrollit",
    "Zgjidhjet aktuale: të shtrenjta, pronësore, të lidhura me një shitës",
]))

# ── SECTION 1 ────────────────────────────────────────────────────────────────
EN.append(("section", "1", "The opensecstack Ecosystem", "Eight platforms. One immutable governance layer."))
AL.append(("section", "1", "Ekosistemi opensecstack", "Tetë platforma. Një shtresë qeverisje e pandryshueshme."))

# ── Slide 5: Ecosystem Overview ──────────────────────────────────────────────
EN.append(("custom_ecosystem", None, None))
AL.append(("custom_ecosystem_al", None, None))

# ── Slide 6: APIGuard ────────────────────────────────────────────────────────
EN.append(("bullets", "APIGuard — API Security Testing", [
    "Parses OpenAPI 3.x, Swagger 2.0, GraphQL, and Postman collections",
    "Full OWASP API Security Top 10 coverage (A1–A10)",
    "  A1: BOLA  ·  A2: Broken Auth  ·  A3: Mass Assignment  ·  A4: Rate Limiting",
    "  A5: Function Auth  ·  A6: Sensitive Flows  ·  A7: SSRF  ·  A8-A10: Advanced",
    "CVSS 3.1 scoring on every finding — prioritise by real risk",
    "Output: SARIF · HTML · JSON · plain-text — native CI/CD integration",
    "React dashboard: scan history, finding trends, regression detection",
    "APIGuard scans itself — the ecosystem's own APIs are continuously tested",
]))
AL.append(("bullets", "APIGuard — Testimi i Sigurisë API", [
    "Analizon OpenAPI 3.x, Swagger 2.0, GraphQL dhe koleksionet Postman",
    "Mbulim i plotë OWASP API Top 10 (A1–A10)",
    "  A1: BOLA  ·  A2: Auth e thyer  ·  A3: Caktim masiv  ·  A4: Kufizim rate",
    "  A5: Auth funksioni  ·  A6: Rrjedha sensitive  ·  A7: SSRF  ·  A8-A10",
    "Vlerësim CVSS 3.1 për çdo gjetje — prioritizim sipas rrezikut real",
    "Dalje: SARIF · HTML · JSON · tekst — integrim nativ CI/CD",
    "Panel React: histori skanimesh, tendenca gjetjesh, zbulim regresioni",
    "APIGuard skanon veten — API-të e ekosistemit testohen vazhdimisht",
]))

# ── Slide 7: NIS2 Compass ───────────────────────────────────────────────────
EN.append(("bullets", "NIS2 Compass — Regulatory Compliance", [
    "Maps all 10 NIS2 Article 21(2) technical measures (A through J)",
    "  (a) Risk analysis  ·  (b) Incident handling  ·  (c) Business continuity",
    "  (d) Supply chain  ·  (e) Network security  ·  (f) Effectiveness",
    "  (g) Cyber hygiene  ·  (h) Cryptography  ·  (i) HR security  ·  (j) Access control",
    "Evidence artefacts — cryptographic hashes stored per control",
    "Assessment workflows with dual-person verifier sign-off",
    "CITADEL webhook: compliance state changes trigger VIGIL advisories",
    "Regulator-ready output: PDF + JSON with full evidence chain",
]))
AL.append(("bullets", "NIS2 Compass — Pajtueshmëria Rregullatore", [
    "Harton të 10 masat teknike të NIS2 neni 21(2) (A deri J)",
    "  (a) Analiza rreziku  ·  (b) Trajtimi incidentit  ·  (c) Vazhdimësia e biznesit",
    "  (d) Zinxhiri furnizimit  ·  (e) Siguria e rrjetit  ·  (f) Efektiviteti",
    "  (g) Higjiena kibernetike  ·  (h) Kriptografia  ·  (i) HR  ·  (j) Kontrolli i aksesit",
    "Artefakte prove — hashes kriptografike ruhen për çdo kontroll",
    "Flukse vlerësimi me nënshkrim verifikuesi me dy persona",
    "Webhook CITADEL: ndryshimet e gjendjes së pajtueshmërisë aktivizojnë këshilla VIGIL",
    "Dalje gati për rregullator: PDF + JSON me zinxhir prove të plotë",
]))

# ── Slide 8: CITADEL ─────────────────────────────────────────────────────────
EN.append(("custom_citadel", None, None))
AL.append(("custom_citadel_al", None, None))

# ── SECTION 2 ────────────────────────────────────────────────────────────────
EN.append(("section", "2", "Five-Layer Defence Model", "Identity · Network · Host · Application · Data"))
AL.append(("section", "2", "Modeli i Mbrojtjes me Pesë Shtresa", "Identiteti · Rrjeti · Hosti · Aplikacioni · Të dhënat"))

# ── Slide 10: Five Layers ────────────────────────────────────────────────────
EN.append(("custom_layers", None, None))
AL.append(("custom_layers_al", None, None))

# ── Slide 11: Layer Details ──────────────────────────────────────────────────
EN.append(("bullets", "Layer-by-Layer: What Stops What", [
    "L1 Identity: SoD enforcement — operator ≠ verifier, cryptographically",
    "  Stolen credential alone is insufficient — the attacker also needs the verifier account",
    "L2 Network: rate limiting, TLS 1.2+, replay protection ±300s window, CORS",
    "  Two-bucket sliding window: (key_id, ts, sig) tuples are one-time use",
    "L3 Host: non-root containers, pinned dependencies, supply-chain checksums",
    "  Kubernetes least-privilege — container escape has minimal blast radius",
    "L4 Application: MARSHAL gates, AUGUR pre-emptive alerts, VIGIL monitoring",
    "  APIGuard continuously scans the ecosystem's own APIs — self-defending",
    "L5 Data: WORM log immutability, TripleHash, Ed25519 anchors, triple backup",
    "  After any recovery: the chain can be cryptographically re-verified",
]))
AL.append(("bullets", "Shtresë pas Shtrese: Çfarë Bllokon Çfarë", [
    "L1 Identiteti: Zbatim NDS — operatori ≠ verifikuesi, kriptografikisht",
    "  Kredencial i vjedhur i vetëm është i pamjaftueshëm — sulmuesit i duhet edhe llogaria verifikuese",
    "L2 Rrjeti: kufizim rate, TLS 1.2+, mbrojtje replay ±300s, CORS",
    "  Dritare rrëshqitëse me dy kova: tuplet (key_id, ts, sig) janë për përdorim njëherësh",
    "L3 Hosti: kontejnerë jo-root, varësi të fiksuara, kontrolle zinxhir furnizimi",
    "  Kubernetes me privilegje minimale — arratisja nga kontejneri ka ndikim minimal",
    "L4 Aplikacioni: porta MARSHAL, alarme parandaluese AUGUR, monitorim VIGIL",
    "  APIGuard skanon vazhdimisht API-të e ekosistemit — vetë-mbrojtëse",
    "L5 Të dhënat: pandryshmëri WORM, TripleHash, ankora Ed25519, kopje rezervë triple",
    "  Pas çdo rikuperimi: zinxhiri mund të ri-verifikohet kriptografikisht",
]))

# ── SECTION 3 ────────────────────────────────────────────────────────────────
EN.append(("section", "3", "The Secure Operating System", "Microkernel · Grid Sandboxes · Ecosystem Integration"))
AL.append(("section", "3", "Sistemi Operativ i Sigurt", "Mikrobërthame · Kuti Rëre Rrjetore · Integrim Ekosistemi"))

# ── Slide 13: Why a New OS ───────────────────────────────────────────────────
EN.append(("bullets", "Why Build a Secure OS from Scratch?", [
    "Modern operating systems were not designed for zero-trust governance",
    "  Monolithic kernels: millions of lines of code in privileged space — huge attack surface",
    "  Retrofitted security: SELinux, AppArmor add policy but not architecture",
    "Our approach: microkernel + capability-based security from the ground up",
    "  Only IPC, memory management, and thread scheduling run in kernel space",
    "  Every driver, filesystem, network stack runs in isolated user-space processes",
    "Grid sandbox model: applications live in cryptographically isolated cells",
    "  Cross-cell communication is governed by CITADEL MARSHAL — logged, auditable",
    "opensecstack runs natively inside the OS — governance is infrastructure, not application",
]))
AL.append(("bullets", "Pse të Ndërtojmë një OS të Sigurt nga E Para?", [
    "Sistemet operative moderne nuk u projektuan për qeverisje zero-trust",
    "  Bërthamat monolitike: miliona rreshta kodi në hapësirë të privilegjuar — sipërfaqe e madhe sulmi",
    "  Siguri e shtuar: SELinux, AppArmor shtojnë politikë por jo arkitekturë",
    "Qasja jonë: mikrobërthame + siguri e bazuar në aftësi nga e para",
    "  Vetëm IPC, menaxhimi i memories dhe planifikimi i thread-eve janë në hapësirën e kernelit",
    "  Çdo shofer, sistem skedarësh, stivë rrjeti funksionon në procese të izoluara të hapësirës së përdoruesit",
    "Modeli sandbox rrjetor: aplikacionet jetojnë në qeliza të izoluara kriptografikisht",
    "  Komunikimi ndër-qelizor qeveriset nga CITADEL MARSHAL — i regjistruar, i auditushëm",
    "opensecstack funksionon në mënyrë vendase brenda OS — qeverisja është infrastrukturë, jo aplikacion",
]))

# ── Slide 14: Microkernel Architecture ──────────────────────────────────────
EN.append(("custom_microkernel", None, None))
AL.append(("custom_microkernel_al", None, None))

# ── Slide 15: Desktop OS Layers ─────────────────────────────────────────────
EN.append(("custom_desktop_os", None, None))
AL.append(("custom_desktop_os_al", None, None))

# ── Slide 16: Mobile OS + MVNO ──────────────────────────────────────────────
EN.append(("custom_mobile_os", None, None))
AL.append(("custom_mobile_os_al", None, None))

# ── Slide 17: Grid Sandboxes ─────────────────────────────────────────────────
EN.append(("bullets", "Grid Sandboxes — Cryptographic Isolation Model", [
    "Each application occupies exactly one sandbox cell in the grid",
    "  Cells are isolated at kernel level — no shared memory, no shared file descriptors",
    "  Cell identity is a cryptographic capability token, not a PID or user ID",
    "Communication between cells requires a MARSHAL-issued channel permit",
    "  Channel permit: signed by CITADEL, time-bounded, audited in WORM log",
    "  No direct inter-process calls — all IPC passes through the governance layer",
    "Sandbox cells are layered: Tier 1 (critical), Tier 2 (trusted), Tier 3 (untrusted)",
    "  Tier 1 ↔ Tier 2: MARSHAL second-hand evaluation (<300ms)",
    "  Tier 2 ↔ Tier 3: MARSHAL with evidence requirement (minute-hand)",
    "VIGIL monitors cell health — anomalous IPC patterns raise AUGUR advisories",
]))
AL.append(("bullets", "Kuti Rëre Rrjetore — Modeli i Izolimit Kriptografik", [
    "Çdo aplikacion zë saktësisht një qelizë sandbox në rrjet",
    "  Qelizat janë të izoluara në nivel kerneli — pa memorie të ndarë, pa deskriptorë skedarësh",
    "  Identiteti i qelizës është një token aftësish kriptografike, jo PID ose ID përdoruesi",
    "Komunikimi ndërmjet qelizave kërkon leje kanali të lëshuar nga MARSHAL",
    "  Leja e kanalit: nënshkruar nga CITADEL, e kufizuar në kohë, e audituar në regjistrin WORM",
    "  Pa thirrje direkte ndër-procese — i gjithë IPC kalon nëpër shtresën e qeverisjes",
    "Qelizat sandbox janë të shtresuara: Tier 1 (kritike), Tier 2 (besuar), Tier 3 (i pabesuar)",
    "  Tier 1 ↔ Tier 2: vlerësim MARSHAL sekondë (<300ms)",
    "  Tier 2 ↔ Tier 3: MARSHAL me kërkesë prove (minutë)",
    "VIGIL monitoron shëndetin e qelizave — modelet anormale IPC ngrenë këshilla AUGUR",
]))

# ── SECTION 4 ────────────────────────────────────────────────────────────────
EN.append(("section", "4", "Blockchain & Triple Hashing", "SHA-256 · SHA-512 · BLAKE3 — Time-Aligned Cryptography"))
AL.append(("section", "4", "Blockchain & Hash-im i Trefishtë", "SHA-256 · SHA-512 · BLAKE3 — Kriptografi e Harmonizuar me Kohën"))

# ── Slide 19: Why Blockchain-Inspired ────────────────────────────────────────
EN.append(("bullets", "Why Blockchain-Inspired Design?", [
    "Traditional audit logs: stored in mutable databases — deleted, edited, fabricated",
    "Public blockchains: decentralised but slow, expensive, complex governance",
    "Our model: cryptographic chain — blockchain's integrity, without its overhead",
    "  Append-only: no UPDATE, no DELETE — enforced at database trigger level",
    "  Chain hashing: changing one entry invalidates every entry that follows",
    "  Anchoring: periodic Ed25519 signatures seal the chain history",
    "The result: a forensically sound, self-verifying audit trail",
    "  Any third party can verify the entire chain with only the genesis hash + public key",
    "  No trusted authority required — mathematical proof is the authority",
]))
AL.append(("bullets", "Pse Dizajn i Frymëzuar nga Blockchain?", [
    "Regjistrat tradicionale të auditimit: ruhen në baza të dhënash të ndryshueshme — fshihen, ndryshohen",
    "Blockchain publike: decentralizuar por e ngadaltë, e shtrenjtë, qeverisje komplekse",
    "Modeli ynë: zinxhir kriptografik — integriteti i blockchain, pa shpenzimet e saj",
    "  Vetëm-shtesë: asnjë UPDATE, asnjë DELETE — zbatuar në nivel trigger baze të dhënash",
    "  Hash zinxhiri: ndryshimi i një hyrjeje shfuqizon çdo hyrje pasardhëse",
    "  Ankorimi: nënshkrime periodike Ed25519 mbyllin historinë e zinxhirit",
    "Rezultati: gjurmë auditimi e shëndoshë nga ana forensike, vetë-verifikuese",
    "  Çdo palë e tretë mund të verifikojë zinxhirin e plotë vetëm me hash gjenezës + çelësin publik",
    "  Nuk kërkohet autoritet i besuar — prova matematike është autoriteti",
]))

# ── Slide 20: SHA-256 ─────────────────────────────────────────────────────────
EN.append(("custom_sha256", None, None))
AL.append(("custom_sha256_al", None, None))

# ── Slide 21: SHA-512 ─────────────────────────────────────────────────────────
EN.append(("custom_sha512", None, None))
AL.append(("custom_sha512_al", None, None))

# ── Slide 22: BLAKE3 ─────────────────────────────────────────────────────────
EN.append(("custom_blake3", None, None))
AL.append(("custom_blake3_al", None, None))

# ── Slide 23: Triple Hash Together ───────────────────────────────────────────
EN.append(("custom_triplehash", None, None))
AL.append(("custom_triplehash_al", None, None))

# ── Slide 24: TDS Alignment ──────────────────────────────────────────────────
EN.append(("custom_tds", None, None))
AL.append(("custom_tds_al", None, None))

# ── SECTION 5 ────────────────────────────────────────────────────────────────
EN.append(("section", "5", "CITADEL Deep Dive", "MARSHAL · VIGIL · AUGUR · WORM"))
AL.append(("section", "5", "CITADEL Thellë", "MARSHAL · VIGIL · AUGUR · WORM"))

# ── Slide 26: MARSHAL ────────────────────────────────────────────────────────
EN.append(("custom_marshal", None, None))
AL.append(("custom_marshal_al", None, None))

# ── Slide 27: VIGIL + DEEP ───────────────────────────────────────────────────
EN.append(("bullets", "VIGIL — Real-Time & Deep Audit", [
    "VIGIL_REALTIME: GREEN · AMBER · RED — refreshed every poll interval",
    "  GREEN: no HARD_STOPs, all mirrors fresh, anchor rotation current",
    "  AMBER: HIGH advisory active, or mirror staleness exceeds threshold",
    "  RED: CRITICAL advisory active, or unresolved HARD_STOP incident",
    "VIGIL_DEEP: daily full-chain audit — runs as hour-hand async job",
    "  Chain integrity: walk every WORM entry, recompute prev→chain hash",
    "  SoD violations: operator_user_id == verifier_user_id in any entry",
    "  Orphaned incidents: HARD_STOP with no subsequent close_incident EXECUTE",
    "  Evidence gaps: EXECUTE outcome with no linked artefacts or change_id",
    "  Deep report: SHA-256 signed, emitted as a WORM entry — tamper-evident",
]))
AL.append(("bullets", "VIGIL — Auditim në Kohë Reale & i Thellë", [
    "VIGIL_REALTIME: JESHILE · PORTOKALLI · E KUQE — rifreskuar çdo interval sondazhi",
    "  JESHILE: asnjë HARD_STOP, të gjitha pasqyrat të freskëta, rrotullimi i ankorës aktual",
    "  PORTOKALLI: këshillë HIGH aktive, ose moshja e pasqyrës tejkalon pragun",
    "  E KUQE: këshillë KRITIKE aktive, ose incident HARD_STOP i pazgjidhur",
    "VIGIL_DEEP: auditim i plotë ditor i zinxhirit — ekzekutohet si detyrë async orare",
    "  Integriteti zinxhirit: ecje nëpër çdo hyrje WORM, rillogaritje hash prev→zinxhir",
    "  Shkelje NDS: operator_user_id == verifier_user_id në çdo hyrje",
    "  Incidente jetimë: HARD_STOP pa close_incident EXECUTE pasues",
    "  Boshllëqe prove: rezultat EXECUTE pa artefakte të lidhura ose change_id",
    "  Raport i thellë: nënshkruar SHA-256, emetuar si hyrje WORM — i dukshëm ndaj manipulimit",
]))

# ── Slide 28: AUGUR ──────────────────────────────────────────────────────────
EN.append(("custom_augur", None, None))
AL.append(("custom_augur_al", None, None))

# ── SECTION 6 ────────────────────────────────────────────────────────────────
EN.append(("section", "6", "MVNO Integration", "Governing Mobile Network Operations with CITADEL"))
AL.append(("section", "6", "Integrimi MVNO", "Qeverisja e Operacioneve të Rrjetit Celular me CITADEL"))

# ── Slide 30: What is MVNO ───────────────────────────────────────────────────
EN.append(("bullets", "What is an MVNO?", [
    "Mobile Virtual Network Operator: provides mobile services without owning spectrum",
    "  Leases radio access from an MNO (e.g. Vodafone, Deutsche Telekom, Orange)",
    "  Full control over: SIM provisioning, billing, data policies, roaming, VoIP",
    "Why MVNO + opensecstack?",
    "  Every SIM activation, data policy change, or roaming agreement is a governance event",
    "  These decisions must be auditable — regulators demand proof of policy enforcement",
    "  MVNO operations touch critical infrastructure — NIS2 Article 21 applies fully",
    "Our model: MVNO operations run inside the opensecstack OS on the microkernel stack",
    "  Network management isolated in a Tier 1 sandbox cell",
    "  All provisioning actions evaluated by CITADEL MARSHAL before execution",
]))
AL.append(("bullets", "Çfarë është MVNO?", [
    "Operator Virtual i Rrjetit Celular: ofron shërbime celulare pa zotëruar spektër",
    "  Merr me qira aksesin radio nga një MNO (p.sh. Vodafone, Deutsche Telekom, Orange)",
    "  Kontroll i plotë mbi: provizionim SIM, faturim, politika të dhënash, roaming, VoIP",
    "Pse MVNO + opensecstack?",
    "  Çdo aktivizim SIM, ndryshim politike të dhënash ose marrëveshje roamingu është ngjarje qeverisëse",
    "  Këto vendime duhet të jenë të auditushme — rregullatorët kërkojnë prova zbatimi politikash",
    "  Operacionet MVNO prekin infrastrukturën kritike — NIS2 neni 21 zbatohet plotësisht",
    "Modeli ynë: operacionet MVNO ekzekutohen brenda OS-it opensecstack në stivën mikrobërthame",
    "  Menaxhimi i rrjetit i izoluar në qelizë sandbox Tier 1",
    "  Të gjitha veprimet e provizionimit vlerësohen nga CITADEL MARSHAL para ekzekutimit",
]))

# ── Slide 31: MVNO Architecture ──────────────────────────────────────────────
EN.append(("custom_mvno", None, None))
AL.append(("custom_mvno_al", None, None))

# ── SECTION 7 ────────────────────────────────────────────────────────────────
EN.append(("section", "7", "EU Alignment & Compliance", "NIS2 · GDPR · Cyber Resilience Act"))
AL.append(("section", "7", "Harmonizimi & Pajtueshmëria BE", "NIS2 · GDPR · Akti i Rezistencës Kibernetike"))

# ── Slide 33: NIS2 Matrix ────────────────────────────────────────────────────
EN.append(("custom_nis2", None, None))
AL.append(("custom_nis2_al", None, None))

# ── Slide 34: GDPR + CRA ─────────────────────────────────────────────────────
EN.append(("bullets", "GDPR, Cyber Resilience Act & EU Data Sovereignty", [
    "GDPR Article 32: appropriate technical measures — TripleHash + WORM log satisfies",
    "  SHA-512 archival hash = long-term integrity of personal data records",
    "  Right to erasure: WORM log stores hashes of data, not data itself",
    "EU Cyber Resilience Act (CRA) — product security requirements",
    "  Secure by design: microkernel + grid sandboxes achieve this by architecture",
    "  Vulnerability reporting: opensecstack VIGIL_DEEP automates continuous assurance",
    "  SBOM: every component is open-source with verifiable supply chain checksums",
    "EU AI Act — governance for AI system actions",
    "  MARSHAL 5-gate engine provides the 'meaningful human oversight' requirement",
    "  Every AI-initiated action: auditable, reversible, with full evidence chain",
]))
AL.append(("bullets", "GDPR, Akti i Rezistencës Kibernetike & Sovraniteti i të Dhënave BE", [
    "GDPR neni 32: masa teknike të përshtatshme — TripleHash + regjistri WORM kënaq",
    "  Hash arkivor SHA-512 = integritet afatgjatë i regjistrimeve të të dhënave personale",
    "  E drejta për fshirje: regjistri WORM ruan hashes të të dhënave, jo vetë të dhënat",
    "Akti i Rezistencës Kibernetike BE (CRA) — kërkesat e sigurisë së produktit",
    "  I sigurt nga dizajni: mikrobërthama + sandbox rrjetore e arrijnë këtë me arkitekturë",
    "  Raportimi i dobësive: VIGIL_DEEP automatizon sigurimin e vazhdueshëm",
    "  SBOM: çdo komponent është me burim të hapur me kontrolle të verifikueshme zinxhir furnizimi",
    "Akti BE i AI — qeverisja për veprimet e sistemit AI",
    "  Motori me 5 porta MARSHAL siguron kërkesën e 'mbikëqyrjes kuptimplotë njerëzore'",
    "  Çdo veprim i iniciuar nga AI: i auditushëm, i kthyeshëm, me zinxhir prove të plotë",
]))

# ── SECTION 8 ────────────────────────────────────────────────────────────────
EN.append(("section", "8", "Roadmap & Call to Action", "Phase 4 · IRFlow · ThreatFlow · v1.0"))
AL.append(("section", "8", "Plani & Thirrja për Veprim", "Faza 4 · IRFlow · ThreatFlow · v1.0"))

# ── Slide 36 (last): CTA ─────────────────────────────────────────────────────
EN.append(("custom_cta", None, None))
AL.append(("custom_cta_al", None, None))


# ══════════════════════════════════════════════════════════════════════════════
# CUSTOM SLIDE RENDERERS
# ══════════════════════════════════════════════════════════════════════════════

def render_ecosystem(prs, lang="EN"):
    s = slide_content(prs, "The opensecstack Ecosystem" if lang=="EN" else "Ekosistemi opensecstack")

    # 8 platforms in 2 rows of 4  ─────────────────────────────────────────────
    # Row 1: core / built
    row1_en = [
        ("APIGuard",      "API Security\nOWASP A1–A10\nCVSS 3.1 · SARIF",         NAVY),
        ("NIS2 Compass",  "Regulatory\nNIS2 Art.21(2)\n10 Measures A–J",           NAVY),
        ("CITADEL",       "Governance Engine\nMARSHAL · WORM\nVIGIL · AUGUR",      NAVY),
        ("IRFlow",        "Incident Response\nPlaybooks · SoD\n[Phase 4]",         DGREY),
    ]
    # Row 2: OS / network / intelligence
    row2_en = [
        ("ThreatFlow",    "Threat Intelligence\nIOC feeds · TTP map\n[Phase 4]",   DGREY),
        ("MVNO Stack",    "Mobile Network Ops\nSIM · eSIM · VoIP\nCITADEL-gated",  NAVY),
        ("Runix\nDesktop", "Microkernel OS\nGrid sandboxes\nHardware-backed",  NAVY),
        ("Runix\nMobile",  "ARM TrustZone\nSIM governance\nMVNO-native",       NAVY),
    ]
    row1_al = [
        ("APIGuard",      "Siguria API\nOWASP A1–A10\nCVSS 3.1 · SARIF",          NAVY),
        ("NIS2 Compass",  "Pajtueshmëri\nNIS2 Neni 21(2)\n10 Masa A–J",            NAVY),
        ("CITADEL",       "Motor Qeverisës\nMARSHAL · WORM\nVIGIL · AUGUR",       NAVY),
        ("IRFlow",        "Reagim Incidenti\nPlanbook · NDS\n[Faza 4]",            DGREY),
    ]
    row2_al = [
        ("ThreatFlow",    "Inteligjencë Kërcënimesh\nFurnizues IOC · TTP\n[Faza 4]", DGREY),
        ("MVNO Stack",    "Operacione Rrjeti Celular\nSIM · eSIM · VoIP\nE kontrolluar CITADEL", NAVY),
        ("Runix\nDesktop", "OS Mikrobërthame\nQeliza sandbox\nMbështetje HW", NAVY),
        ("Runix\nMobile",  "ARM TrustZone\nQeverisja SIM\nMVNO-nativ",        NAVY),
    ]

    row1 = row1_en if lang == "EN" else row1_al
    row2 = row2_en if lang == "EN" else row2_al

    BW = Inches(2.9)   # box width
    BH = Inches(1.75)  # box height
    GAP = Inches(0.51) # gap between boxes (room for arrow)
    LEFT = Inches(0.32)
    R1Y = Inches(1.3)
    R2Y = Inches(3.25)

    for i, (lbl, val, clr) in enumerate(row1):
        x = LEFT + i * (BW + GAP)
        box(s, lbl, val, x, R1Y, w=BW, h=BH, bg=clr, ls=11, vs=13)
        if i < 3:
            rect(s, x + BW, R1Y + BH / 2 - Inches(0.04), GAP, Inches(0.08), TEAL)

    for i, (lbl, val, clr) in enumerate(row2):
        x = LEFT + i * (BW + GAP)
        box(s, lbl, val, x, R2Y, w=BW, h=BH, bg=clr, ls=11, vs=13)
        if i < 3:
            rect(s, x + BW, R2Y + BH / 2 - Inches(0.04), GAP, Inches(0.08), TEAL)

    note = ("All platforms share one CITADEL WORM chain · HMAC-SHA256 connector auth · Apache 2.0 · SDK: Go · Python · Rust"
            if lang=="EN" else
            "Të gjitha platformat ndajnë një zinxhir WORM CITADEL · Auth HMAC-SHA256 · Apache 2.0 · SDK: Go · Python · Rust")
    txt(s, note, Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.45),
        size=13, italic=True, color=DGREY, align=PP_ALIGN.CENTER)

    stats_en = [("8", "Platforms"), ("3", "SDK Languages"), ("5", "MARSHAL Gates"),
                ("9", "AUGUR Rules"), ("3", "Hash Algorithms")]
    stats_al = [("8", "Platforma"), ("3", "Gjuhë SDK"), ("5", "Porta MARSHAL"),
                ("9", "Rregulla AUGUR"), ("3", "Algoritme Hash")]
    stats = stats_en if lang == "EN" else stats_al
    for i, (val, lbl) in enumerate(stats):
        bx = Inches(0.28 + i * 2.6)
        rect(s, bx, Inches(5.8), Inches(2.45), Inches(1.55), TEAL)
        txt(s, val, bx, Inches(5.85), Inches(2.45), Inches(0.72),
            size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(s, lbl, bx, Inches(6.57), Inches(2.45), Inches(0.72),
            size=12, color=NAVY, align=PP_ALIGN.CENTER)

def render_citadel(prs, lang="EN"):
    t = "CITADEL — The Governance Engine" if lang=="EN" else "CITADEL — Motori i Qeverisjes"
    s = slide_content(prs, t)
    comps_en = [("MARSHAL","5-Gate Decision\nEngine"), ("WORM Log","Append-Only\nAudit Chain"),
                ("VIGIL","Real-Time\nMonitor"), ("AUGUR","9 Advisory\nRules"),
                ("Anchor","Ed25519\nKey Rotation")]
    comps_al = [("MARSHAL","Motor Vendimmarrës\n5 Portash"), ("Regjistri WORM","Zinxhir\nAuditimi"),
                ("VIGIL","Monitor\nKohë Reale"), ("AUGUR","9 Rregulla\nKëshillimi"),
                ("Ankora","Rrotullim Çelësi\nEd25519")]
    comps = comps_en if lang=="EN" else comps_al
    colors = [GREEN, BLUE, AMBER, RED, PURPLE]
    for i, ((lbl, val), clr) in enumerate(zip(comps, colors)):
        x = Inches(0.25 + i * 2.6)
        box(s, lbl, val, x, Inches(1.35), w=Inches(2.45), h=Inches(1.8), bg=clr)
    arrow = ("⟵── All components write signed events to the WORM log ──⟶"
             if lang=="EN" else "⟵── Të gjithë komponentët shkruajnë ngjarje të nënshkruara në regjistrin WORM ──⟶")
    txt(s, arrow, Inches(0.3), Inches(3.35), Inches(12.7), Inches(0.5),
        size=14, italic=True, color=TEAL, align=PP_ALIGN.CENTER)
    auth_en = [
        "Connector auth: HMAC-SHA256 — X-CITADEL-KEY · X-CITADEL-TS · X-CITADEL-SIG",
        "Admin auth: HS256 JWT — group_sig_admin · group_sig_auditor roles",
        "Two-bucket sliding window: (key_id, ts, sig) nonces are one-time only",
        "All decisions immutable: EXECUTE · REFUSE · HARD_STOP — WORM-logged forever",
    ]
    auth_al = [
        "Auth konektori: HMAC-SHA256 — X-CITADEL-KEY · X-CITADEL-TS · X-CITADEL-SIG",
        "Auth admin: HS256 JWT — rolet group_sig_admin · group_sig_auditor",
        "Dritare rrëshqitëse dy-kova: nonce-t (key_id, ts, sig) janë njëherësh",
        "Të gjitha vendimet të pandryshueshme: EXECUTE · REFUSE · HARD_STOP — WORM-logged",
    ]
    auth = auth_en if lang=="EN" else auth_al
    y = Inches(4.0)
    for line in auth:
        txt(s, "  ▶  " + line, Inches(0.4), y, Inches(12.5), Inches(0.5),
            size=16, color=DGREY, bold=False)
        y += Inches(0.52)

def render_layers(prs, lang="EN"):
    t = "Five-Layer Defence — Concentric Security" if lang=="EN" else "Mbrojtja me Pesë Shtresa — Siguri Koncentrike"
    s = slide_content(prs, t)
    layers_en = [
        ("L5  DATA FENCE",        "WORM log · TripleHash · Ed25519 anchors · Triple backup · DR",   TEAL,   Inches(0.3)),
        ("L4  APPLICATION FENCE", "APIGuard · MARSHAL gates · VIGIL · AUGUR · Input validation",    BLUE,   Inches(0.7)),
        ("L3  HOST FENCE",        "Kubernetes · Least-privilege containers · Pinned dependencies",   PURPLE, Inches(1.1)),
        ("L2  NETWORK FENCE",     "Rate limiting · TLS 1.2+ · Replay protection · CORS · Isolation", GREEN,  Inches(1.5)),
        ("L1  IDENTITY FENCE",    "JWT · API keys · HMAC-SHA256 · SoD enforcement · Role groups",    AMBER,  Inches(1.9)),
    ]
    layers_al = [
        ("L5  GARDHI I TË DHËNAVE", "WORM · TripleHash · ankora Ed25519 · Kopje rezervë triple · DR",     TEAL,   Inches(0.3)),
        ("L4  GARDHI APLIKACIONIT", "APIGuard · Porta MARSHAL · VIGIL · AUGUR · Validim hyrjeje",          BLUE,   Inches(0.7)),
        ("L3  GARDHI I HOSTIT",     "Kubernetes · Kontejnerë privilegje-minimale · Varësi të fiksuara",     PURPLE, Inches(1.1)),
        ("L2  GARDHI I RRJETIT",    "Kufizim rate · TLS 1.2+ · Mbrojtje replay · CORS · Izolim",           GREEN,  Inches(1.5)),
        ("L1  GARDHI IDENTITETIT",  "JWT · Çelësa API · HMAC-SHA256 · Zbatim NDS · Grupet e roleve",       AMBER,  Inches(1.9)),
    ]
    layers = layers_en if lang == "EN" else layers_al
    for lbl, desc, clr, x in layers:
        w = W - x * 2
        rect(s, x, Inches(1.35) + (layers.index((lbl,desc,clr,x)) * Inches(1.0)),
             w, Inches(0.92), clr)
        idx = layers.index((lbl, desc, clr, x))
        y2 = Inches(1.35) + idx * Inches(1.0)
        txt(s, lbl,  x + Inches(0.12), y2 + Inches(0.07), Inches(3.5), Inches(0.4),
            size=15, bold=True, color=WHITE)
        txt(s, desc, x + Inches(3.7),  y2 + Inches(0.07), w - Inches(4.0), Inches(0.78),
            size=14, color=WHITE)

def render_microkernel(prs, lang="EN"):
    t = "Microkernel Architecture — Minimal Trusted Core" if lang=="EN" else "Arkitektura Mikrobërthame — Bërthamë Minimale e Besuar"
    s = slide_content(prs, t)
    kern_en  = [
        ("KERNEL SPACE\n(Trusted Core)", "IPC · Memory Management\nThread Scheduling only",             NAVY,   Inches(4.2), Inches(1.35), Inches(4.9), Inches(1.3)),
    ]
    kern_al  = [
        ("HAPËSIRA E KERNELIT\n(Bërthamë e Besuar)", "IPC · Menaxhim Memorje\nPlanifikim Thread-esh",   NAVY,   Inches(4.2), Inches(1.35), Inches(4.9), Inches(1.3)),
    ]
    kern = kern_en if lang=="EN" else kern_al
    for lbl, desc, clr, x, y, w2, h2 in kern:
        box(s, lbl, desc, x, y, w=w2, h=h2, bg=clr, lc=TEAL, vc=WHITE, ls=14, vs=15)

    user_en = [
        ("Drivers",       "GPU · NIC · Storage\nin user space",       BLUE,   Inches(0.3),  Inches(3.0), Inches(2.7), Inches(1.5)),
        ("File Systems",  "ext4 · NTFS · tmpfs\nisolated processes",  GREEN,  Inches(3.2),  Inches(3.0), Inches(2.7), Inches(1.5)),
        ("Network Stack", "TCP/IP · TLS layer\nin sandbox cell",      PURPLE, Inches(6.1),  Inches(3.0), Inches(2.7), Inches(1.5)),
        ("CITADEL",       "Governance layer\nnative kernel service",  TEAL,   Inches(9.0),  Inches(3.0), Inches(2.7), Inches(1.5)),
        ("App Runtime",   "Grid-sandboxed\nuser applications",        AMBER,  Inches(6.1),  Inches(4.75),Inches(5.6), Inches(1.1)),
    ]
    user_al = [
        ("Shoferë",       "GPU · NIC · Ruajtja\nnë hapësirën user",   BLUE,   Inches(0.3),  Inches(3.0), Inches(2.7), Inches(1.5)),
        ("Sistemet Skedarësh", "ext4 · NTFS · tmpfs\nprocese të izoluara", GREEN, Inches(3.2), Inches(3.0), Inches(2.7), Inches(1.5)),
        ("Stiva Rrjetit", "TCP/IP · shtresë TLS\nnë qelizë sandbox",  PURPLE, Inches(6.1),  Inches(3.0), Inches(2.7), Inches(1.5)),
        ("CITADEL",       "Shtresa qeverisëse\nshërbim nativ kernel", TEAL,   Inches(9.0),  Inches(3.0), Inches(2.7), Inches(1.5)),
        ("Runtime App",   "Aplikacione\nnë sandbox rrjetor",          AMBER,  Inches(6.1),  Inches(4.75),Inches(5.6), Inches(1.1)),
    ]
    user = user_en if lang=="EN" else user_al
    lbl_u = "USER SPACE — All services isolated, communicate via IPC only" if lang=="EN" else "HAPËSIRA USER — Të gjitha shërbimet të izoluara, komunikojnë vetëm nëpërmjet IPC"
    txt(s, lbl_u, Inches(0.3), Inches(2.75), Inches(12.7), Inches(0.3),
        size=12, italic=True, color=DGREY)
    for lbl, desc, clr, x, y, w2, h2 in user:
        box(s, lbl, desc, x, y, w=w2, h=h2, bg=clr)

    note = ("vs Monolithic: millions of lines in kernel space — single CVE can own the whole machine"
            if lang=="EN" else "vs Monolitike: miliona rreshta në hapësirën kernel — një CVE mund të marrë kontrollin e të gjithë makinës")
    txt(s, note, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.45),
        size=14, italic=True, color=RED, align=PP_ALIGN.CENTER)

def render_desktop_os(prs, lang="EN"):
    t = "Desktop OS — 6 Security Layers" if lang=="EN" else "OS Desktop — 6 Shtresa Sigurie"
    s = slide_content(prs, t)
    layers_en = [
        ("L6  USER INTERFACE",    "App shell · Dashboard · Browser — sandboxed cells",                   RGBColor(0x22,0x66,0xAA)),
        ("L5  APP RUNTIME",       "Language runtimes · Containers — each in grid cell, MARSHAL-gated",   BLUE),
        ("L4  CITADEL LAYER",     "MARSHAL · VIGIL · AUGUR — governance native to the kernel",           TEAL),
        ("L3  OS SERVICES",       "Filesystems · Network · Audio — isolated user-space processes",       PURPLE),
        ("L2  MICROKERNEL",       "IPC · Memory · Scheduling — <20k lines, formally verifiable",         GREEN),
        ("L1  HARDWARE / HW SEC", "Secure Boot · TPM 2.0 · HSM — Ed25519 keys in hardware",              AMBER),
    ]
    layers_al = [
        ("L6  NDËRFAQJA PËRDORUESI", "Shell aplikacionesh · Panel · Shfletues — qeliza sandbox",         RGBColor(0x22,0x66,0xAA)),
        ("L5  RUNTIME APLIKACIONI",  "Runtime gjuhësh · Kontejnerë — secili në qelizë rrjetore",        BLUE),
        ("L4  SHTRESA CITADEL",      "MARSHAL · VIGIL · AUGUR — qeverisje native e kernelit",           TEAL),
        ("L3  SHËRBIME OS",          "Sistemet skedarësh · Rrjeti · Audio — procese të izoluara",       PURPLE),
        ("L2  MIKROBËRTHAMË",        "IPC · Memorie · Planifikim — <20k rreshta, e verifikueshme",      GREEN),
        ("L1  HARDWARE / SIGURI HW", "Nisje e sigurt · TPM 2.0 · HSM — çelësat Ed25519 në hardware",   AMBER),
    ]
    layers = layers_en if lang=="EN" else layers_al
    for i, (lbl, desc, clr) in enumerate(layers):
        y = Inches(1.35) + i * Inches(0.97)
        rect(s, Inches(0.3), y, Inches(12.7), Inches(0.9), clr)
        txt(s, lbl,  Inches(0.45), y + Inches(0.08), Inches(3.8), Inches(0.4),
            size=14, bold=True, color=WHITE)
        txt(s, desc, Inches(4.3),  y + Inches(0.08), Inches(8.5), Inches(0.78),
            size=14, color=WHITE)

def render_mobile_os(prs, lang="EN"):
    t = "Mobile OS — Layers + MVNO Integration" if lang=="EN" else "OS Celular — Shtresat + Integrimi MVNO"
    s = slide_content(prs, t)
    layers_en = [
        ("L6  USER INTERFACE",     "Secure dialer · Messaging · App launcher",     RGBColor(0x22,0x66,0xAA)),
        ("L5  APP RUNTIME",        "App sandbox cells — each MARSHAL-gated",        BLUE),
        ("L4  CITADEL LAYER",      "MARSHAL · VIGIL — SIM & roaming governance",    TEAL),
        ("L3  MVNO STACK",         "SIM provisioning · Data policies · VoIP · eSIM",GREEN),
        ("L2  RADIO ABSTRACTION",  "RIL (Radio Interface Layer) — isolated process", PURPLE),
        ("L1  MICROKERNEL + HW",   "ARM TrustZone · Secure Enclave · TPM",          AMBER),
    ]
    layers_al = [
        ("L6  NDËRFAQJA PËRDORUESI","Telefonues i sigurt · Mesazhe · Lëshues aplikacionesh", RGBColor(0x22,0x66,0xAA)),
        ("L5  RUNTIME APLIKACIONI", "Qeliza sandbox aplikacionesh — secili i kontrolluar nga MARSHAL", BLUE),
        ("L4  SHTRESA CITADEL",     "MARSHAL · VIGIL — qeverisja SIM & roamingu",           TEAL),
        ("L3  STIVA MVNO",          "Provizionim SIM · Politika të dhënash · VoIP · eSIM",   GREEN),
        ("L2  ABSTRAKSION RADIO",   "RIL (Shtresa Radio Interface) — proces i izoluar",      PURPLE),
        ("L1  MIKROBËRTHAMË + HW",  "ARM TrustZone · Enklavë e sigurt · TPM",               AMBER),
    ]
    layers = layers_en if lang=="EN" else layers_al
    for i, (lbl, desc, clr) in enumerate(layers):
        y = Inches(1.35) + i * Inches(0.97)
        rect(s, Inches(0.3), y, Inches(8.5), Inches(0.9), clr)
        txt(s, lbl,  Inches(0.45), y + Inches(0.08), Inches(3.0), Inches(0.4),
            size=13, bold=True, color=WHITE)
        txt(s, desc, Inches(3.5),  y + Inches(0.08), Inches(5.1), Inches(0.78),
            size=13, color=WHITE)

    mvno_title = "MVNO Governance" if lang=="EN" else "Qeverisja MVNO"
    txt(s, mvno_title, Inches(9.1), Inches(1.35), Inches(4.1), Inches(0.5),
        size=16, bold=True, color=NAVY)
    mvno_en = ["SIM activation →\nMARSHAL EXECUTE", "Data policy →\nWORM logged", "Roaming →\nSoD enforced",
               "eSIM provision →\nAudit trail", "VoIP call →\nEncrypted + hashed"]
    mvno_al = ["Aktiv. SIM →\nMARSHAL EXECUTE", "Politika →\nWORM i regj.", "Roaming →\nNDS zbatuar",
               "Proviz. eSIM →\nGjurmë auditimi", "Thirrje VoIP →\nI enkriptuar + hash"]
    mvno = mvno_en if lang=="EN" else mvno_al
    for i, m in enumerate(mvno):
        y = Inches(1.9) + i * Inches(1.05)
        rect(s, Inches(9.1), y, Inches(4.0), Inches(0.95), NAVY)
        txt(s, m, Inches(9.2), y + Inches(0.05), Inches(3.8), Inches(0.85),
            size=13, color=WHITE)

def render_hash(prs, algo, lang="EN"):
    names = {"sha256": ("SHA-256 — WORM Chain Integrity",
                        "SHA-256 — Integriteti i Zinxhirit WORM"),
             "sha512": ("SHA-512 — Archival & Anchor Signing",
                        "SHA-512 — Arkivim & Nënshkrim Ankore"),
             "blake3": ("BLAKE3 — Real-Time Speed & Security",
                        "BLAKE3 — Shpejtësi & Siguri Kohë Reale")}
    t = names[algo][0] if lang=="EN" else names[algo][1]
    s = slide_content(prs, t)

    specs = {
        "sha256": {
            "EN": [
                ("Output size",    "256 bits / 32 bytes"),
                ("TDS tier",       "MINUTE HAND  (300ms – 30s)"),
                ("Used for",       "WORM chain hashing — chain_hash = SHA-256(prev || payload)"),
                ("Property",       "Collision resistant — changing any byte changes the entire digest"),
                ("Standardised",   "NIST FIPS 180-4 — accepted in all EU regulatory frameworks"),
                ("Where used",     "Every WORM entry · Connector auth body hash · VIGIL_DEEP report hash"),
            ],
            "AL": [
                ("Madhësia e daljes","256 bit / 32 bajt"),
                ("Niveli TDS",       "MINUTË  (300ms – 30s)"),
                ("Përdoret për",     "Hash zinxhiri WORM — chain_hash = SHA-256(prev || payload)"),
                ("Veti",             "Rezistent ndaj kolizioneve — ndryshimi i çdo bajti ndryshon të gjithë digjestinin"),
                ("I standardizuar",  "NIST FIPS 180-4 — pranuar në të gjitha kornizat rregullatore BE"),
                ("Ku përdoret",      "Çdo hyrje WORM · Hash trup auth konektori · Hash raporti VIGIL_DEEP"),
            ]
        },
        "sha512": {
            "EN": [
                ("Output size",    "512 bits / 64 bytes"),
                ("TDS tier",       "HOUR HAND  (>30s) — archival and anchor signing"),
                ("Used for",       "Long-term integrity of archived evidence and Ed25519 anchor payloads"),
                ("Property",       "128-bit security margin — quantum-resistant for 40+ year retention"),
                ("Standardised",   "NIST FIPS 180-4 — required for EU long-term evidence under eIDAS"),
                ("Where used",     "TripleHash composite · VIGIL_DEEP archival reports · Signed chain exports"),
            ],
            "AL": [
                ("Madhësia e daljes","512 bit / 64 bajt"),
                ("Niveli TDS",       "ORË  (>30s) — arkivim dhe nënshkrim ankore"),
                ("Përdoret për",     "Integritet afatgjatë i provave të arkivuara dhe ngarkesave ankore Ed25519"),
                ("Veti",             "Marzh sigurie 128-bit — rezistent ndaj kuantumit për ruajtje 40+ vjeçare"),
                ("I standardizuar",  "NIST FIPS 180-4 — kërkohet për provat afatgjata BE sipas eIDAS"),
                ("Ku përdoret",      "Kompoziti TripleHash · Raporte arkivale VIGIL_DEEP · Eksporte zinxhiri"),
            ]
        },
        "blake3": {
            "EN": [
                ("Output size",    "256 bits / 32 bytes (configurable up to 128 bytes)"),
                ("TDS tier",       "SECOND HAND  (<300ms) — per-request real-time"),
                ("Used for",       "Real-time integrity checks — WORM payload verification at ingestion"),
                ("Property",       "12× faster than SHA-256 on modern hardware — SIMD parallelism"),
                ("vs SHA-3",       "Same security level, vastly higher throughput — ideal for streaming data"),
                ("Where used",     "TripleHash composite · TDS scanner latency checks · Per-request hashing"),
            ],
            "AL": [
                ("Madhësia e daljes","256 bit / 32 bajt (e konfigurueshme deri 128 bajt)"),
                ("Niveli TDS",       "SEKONDË  (<300ms) — kohë reale për kërkesë"),
                ("Përdoret për",     "Kontroll integriteti kohë reale — verifikim ngarkese WORM në ingestion"),
                ("Veti",             "12× më i shpejtë se SHA-256 në hardware modern — paralelizëm SIMD"),
                ("vs SHA-3",         "I njëjti nivel sigurie, xhiro shumë më e lartë — ideal për të dhëna transmetimi"),
                ("Ku përdoret",      "Kompoziti TripleHash · Kontrolle latence tds-scanner · Hash për kërkesë"),
            ]
        }
    }

    clrs = {"sha256": BLUE, "sha512": PURPLE, "blake3": GREEN}
    clr = clrs[algo]
    spec = specs[algo][lang]
    y = Inches(1.5)
    for lbl, val in spec:
        rect(s, Inches(0.4), y, Inches(2.8), Inches(0.72), clr)
        txt(s, lbl, Inches(0.5), y + Inches(0.07), Inches(2.6), Inches(0.6),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(s, Inches(3.3), y, Inches(9.6), Inches(0.72), LGREY)
        txt(s, val, Inches(3.45), y + Inches(0.07), Inches(9.3), Inches(0.6),
            size=16, color=NAVY)
        y += Inches(0.83)

    algo_colors = {"sha256": BLUE, "sha512": PURPLE, "blake3": GREEN}
    note_en = {"sha256": "SHA-256 is the chain's backbone — every WORM entry depends on every previous hash",
               "sha512": "SHA-512 makes the archive unforgeable across decades — even by quantum computers",
               "blake3": "BLAKE3 is the speed layer — real-time verification without latency cost"}
    note_al = {"sha256": "SHA-256 është shtylla kurrizore e zinxhirit — çdo hyrje WORM varet nga çdo hash i mëparshëm",
               "sha512": "SHA-512 e bën arkivin të pafalshëm nëpër dekada — edhe nga kompjuterët kuantikë",
               "blake3": "BLAKE3 është shtresa e shpejtësisë — verifikim kohë reale pa kosto latence"}
    note = note_en[algo] if lang=="EN" else note_al[algo]
    txt(s, note, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.45),
        size=14, bold=True, color=clr, align=PP_ALIGN.CENTER)

def render_triplehash(prs, lang="EN"):
    t = "TripleHash — 128-byte Composite Digest" if lang=="EN" else "TripleHash — Digjest Kompozit 128 Bajt"
    s = slide_content(prs, t)
    intro = ("TripleHash combines all three algorithms into a single 128-byte composite — all three must be broken simultaneously to forge any record."
             if lang=="EN" else
             "TripleHash kombinon të tre algoritmet në një digjest kompozit 128 bajt — të tre duhet të thyhen njëkohësisht për të falsifikuar çdo regjistrim.")
    txt(s, intro, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.7),
        size=16, color=DGREY, italic=True)
    hashes = [
        ("BLAKE3",  "32 bytes", "Real-time\nSecond hand",   GREEN,  Inches(0.4)),
        ("+", "SHA-256", "32 bytes",                         WHITE,  Inches(3.8)),
        ("SHA-256", "32 bytes", "WORM chain\nMinute hand",   BLUE,   Inches(4.5)),
        ("+", "SHA-512", "64 bytes",                         WHITE,  Inches(7.9)),
        ("SHA-512", "64 bytes", "Archival\nHour hand",       PURPLE, Inches(8.6)),
        ("=", "128 bytes composite", "",                     WHITE,  Inches(11.0)),
    ]
    for i, item in enumerate(hashes):
        if item[0] in ("+", "="):
            txt(s, item[0], item[4], Inches(2.5), Inches(0.6), Inches(1.0),
                size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
            continue
        lbl, sz, tier, clr, x = item
        rect(s, x, Inches(2.2), Inches(3.3), Inches(2.0), clr)
        txt(s, lbl,  x + Inches(0.1), Inches(2.28), Inches(3.1), Inches(0.65),
            size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sz,   x + Inches(0.1), Inches(2.93), Inches(3.1), Inches(0.45),
            size=16, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, tier, x + Inches(0.1), Inches(3.4), Inches(3.1), Inches(0.75),
            size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    result_en = [
        "Any single hash broken → record still protected by remaining two algorithms",
        "SHA-256 + SHA-512 are FIPS-approved → EU regulatory compliance guaranteed",
        "BLAKE3 provides sub-millisecond integrity checking at ingestion time",
        "128-byte digest stored in WORM log — verifiable offline, forever",
    ]
    result_al = [
        "Nëse thyhet ndonjë hash i vetëm → regjistrimi mbetet i mbrojtur nga dy algoritmet e tjera",
        "SHA-256 + SHA-512 janë aprovuar FIPS → pajtueshmëria rregullatore BE e garantuar",
        "BLAKE3 ofron kontroll integriteti nën-milisekondësh në kohën e injektimit",
        "Digjest 128 bajt ruhet në regjistrin WORM — i verifikueshëm jashtë linje, përgjithmonë",
    ]
    result = result_en if lang=="EN" else result_al
    y = Inches(4.6)
    for line in result:
        txt(s, "  ▶  " + line, Inches(0.4), y, Inches(12.5), Inches(0.48),
            size=16, color=NAVY)
        y += Inches(0.5)

def render_tds(prs, lang="EN"):
    t = "Time Dimension Segmentation — Hash + Latency Alignment" if lang=="EN" else "Segmentimi i Dimensionit Kohor — Harmonizimi Hash + Latencë"
    s = slide_content(prs, t)
    intro = ("TDS maps every operation to a latency tier — the hash algorithm chosen must match the tier."
             if lang=="EN" else "TDS harton çdo operacion në një nivel latence — algoritmi hash i zgjedhur duhet të përputhet me nivelin.")
    txt(s, intro, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.5), size=16, color=DGREY, italic=True)
    tds_en = [
        ("SECOND HAND",  "< 300ms",     "BLAKE3",  "Per-request WORM ingestion\nSandbox IPC verification\nMARSHAL pre-check",  GREEN),
        ("MINUTE HAND",  "300ms – 30s", "SHA-256", "WORM chain hashing\nVigilStatus refresh\nConnector key lookup",            BLUE),
        ("HOUR HAND",    "> 30s",       "SHA-512", "VIGIL_DEEP full audit\nArchival signing\nTripleHash composite",             PURPLE),
    ]
    tds_al = [
        ("SEKONDARI",  "< 300ms",     "BLAKE3",  "WORM ingestion për kërkesë\nVerifikim IPC Sandbox\nKontroll paraprak MARSHAL", GREEN),
        ("MINUTAR",    "300ms – 30s", "SHA-256", "Hash zinxhiri WORM\nRifreskim VigilStatus\nKërkesë çelës konektori",           BLUE),
        ("ORAR",       "> 30s",       "SHA-512", "Auditim i plotë VIGIL_DEEP\nNënshkrim arkivor\nKompoziti TripleHash",          PURPLE),
    ]
    tds = tds_en if lang=="EN" else tds_al
    for i, (tier, bound, algo, ops, clr) in enumerate(tds):
        x = Inches(0.4 + i * 4.3)
        rect(s, x, Inches(2.0), Inches(4.1), Inches(4.2), clr)
        txt(s, tier,  x, Inches(2.08), Inches(4.1), Inches(0.55), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, bound, x, Inches(2.62), Inches(4.1), Inches(0.4),  size=14, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
        rect(s, x + Inches(0.9), Inches(3.1), Inches(2.3), Inches(0.55), NAVY)
        txt(s, algo,  x + Inches(0.9), Inches(3.13), Inches(2.3), Inches(0.5), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, ops,   x + Inches(0.1), Inches(3.82), Inches(3.9), Inches(2.2), size=14, color=WHITE)

    note = ("tds-scanner tool measures actual latencies in CI — any second-hand operation exceeding 300ms fails the build."
            if lang=="EN" else "Mjeti tds-scanner mat latencitë aktuale në CI — çdo operacion sekondari që kalon 300ms dështon ndërtimin.")
    txt(s, note, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5),
        size=14, italic=True, color=AMBER, align=PP_ALIGN.CENTER)

def render_marshal(prs, lang="EN"):
    t = "MARSHAL — 5-Gate Decision Engine" if lang=="EN" else "MARSHAL — Motor Vendimmarrës me 5 Porta"
    s = slide_content(prs, t)
    gates_en = [
        ("G1", "AUTHORITY",     "SoD:\nOperator ≠ Verifier\nRole validation",      GREEN),
        ("G2", "SCOPE",         "Project whitelist\nFreeze detection\nAction type", BLUE),
        ("G3", "DETERMINISM",   "Change/Incident ID\nTrace to ticket\nNo orphans",  AMBER),
        ("G4", "EVIDENCE",      "Artefact hashes\nChange ID link\nAudit trail",     PURPLE),
        ("G5", "SCHEMA",        "Kerkese 2.0\nJSON schema\nField validity",         RED),
    ]
    gates_al = [
        ("G1", "AUTORITETI",   "NDS:\nOperator ≠ Verifikues\nVaildim roli",         GREEN),
        ("G2", "FUSHËVEPRIMI", "Lista e lejuar e projektit\nZbulim ngrirjeje",       BLUE),
        ("G3", "DETERMINIZMI", "ID Ndryshimi/Incidenti\nGjurmë deri te bileta",     AMBER),
        ("G4", "PROVA",        "Hashes artefaktesh\nLidhja ID ndryshimi",            PURPLE),
        ("G5", "SKEMA",        "Kerkese 2.0\nSkemë JSON\nVlefshmëri fushash",       RED),
    ]
    gates = gates_en if lang=="EN" else gates_al
    for i, (num, name, desc, clr) in enumerate(gates):
        x = Inches(0.25 + i * 2.6)
        rect(s, x, Inches(1.35), Inches(2.45), Inches(3.7), clr)
        txt(s, num,  x, Inches(1.42), Inches(2.45), Inches(0.45), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, name, x, Inches(1.87), Inches(2.45), Inches(0.6),  size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, x, Inches(2.5),  Inches(2.45), Inches(2.2),  size=14, color=WHITE, align=PP_ALIGN.CENTER)

    outcomes_en = "Outcome: EXECUTE  ·  REFUSE  ·  HARD_STOP — every decision WORM-logged, chain-hashed, permanent"
    outcomes_al = "Rezultati: EXECUTE · REFUSE · HARD_STOP — çdo vendim i regjistruar WORM, hash zinxhiri, i përhershëm"
    txt(s, outcomes_en if lang=="EN" else outcomes_al,
        Inches(0.3), Inches(5.28), Inches(12.7), Inches(0.5),
        size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    note_en = "HARD_STOP is never suppressed — even on internal error, the record is written before the exception propagates"
    note_al = "HARD_STOP nuk supreson kurrë — edhe kur ndodh gabim i brendshëm, regjistrimi shkruhet para se të përpunohet përjashtimi"
    txt(s, note_en if lang=="EN" else note_al,
        Inches(0.3), Inches(5.92), Inches(12.7), Inches(0.55),
        size=14, italic=True, color=DGREY, align=PP_ALIGN.CENTER)

def render_augur(prs, lang="EN"):
    t = "AUGUR — 9 Predictive Advisory Rules" if lang=="EN" else "AUGUR — 9 Rregulla Këshillimi Parashikues"
    s = slide_content(prs, t)
    advisories_en = [
        ("AUG-001","MEDIUM",  "Mandate expiry approaching — renew before MARSHAL Gate 1 refuses"),
        ("AUG-002","HIGH",    "Approval threshold exceeded — additional sign-off required"),
        ("AUG-003","MEDIUM",  "Vendor contract expiry — supply chain risk elevation"),
        ("AUG-004","HIGH",    "Freeze active — MARSHAL Gate 2 will REFUSE non-emergency actions"),
        ("AUG-005","MEDIUM",  "Mirror data inconsistency — cross-check ERP before proceeding"),
        ("AUG-006","MEDIUM",  "Unusual action frequency — possible anomalous automation detected"),
        ("AUG-007","HIGH",    "Mirror stale > threshold — advisory data may be outdated"),
        ("AUG-008","CRITICAL","Unresolved HARD_STOP — all non-emergency actions at elevated risk"),
        ("AUG-009","HIGH",    "Chain anchor overdue — rotation required, timestamp trust weakens"),
    ]
    advisories_al = [
        ("AUG-001","MEDIUM",   "Mandati po skadon — rinovoni para se Porta 1 MARSHAL të refuzojë"),
        ("AUG-002","HIGH",     "Pragu i miratimit tejkaluar — kërkohet nënshkrim shtesë"),
        ("AUG-003","MEDIUM",   "Kontrata e furnizuesit po skadon — ngritje rreziku zinxhiri furnizimit"),
        ("AUG-004","HIGH",     "Ngrirja aktive — Porta 2 MARSHAL do të REFUZOJË veprimet jo-emergjente"),
        ("AUG-005","MEDIUM",   "Mospërputhje të dhënash pasqyre — kontrolloni ERP para vazhdimit"),
        ("AUG-006","MEDIUM",   "Frekuencë e pazakontë veprimesh — zbuluar automatizim anormal"),
        ("AUG-007","HIGH",     "Pasqyra e vjetër > pragu — të dhënat këshilluese mund të jenë të vjetra"),
        ("AUG-008","CRITICAL", "HARD_STOP i pazgjidhur — të gjitha veprimet jo-emergjente me rrezik të ngritur"),
        ("AUG-009","HIGH",     "Ankora e zinxhirit e vonuar — kërkohet rrotullim, besimi i vulës kohore dobësohet"),
    ]
    advisories = advisories_en if lang=="EN" else advisories_al
    sev_c = {"CRITICAL": RED, "HIGH": AMBER, "MEDIUM": TEAL}
    y = Inches(1.4)
    for code, sev, desc in advisories:
        clr = sev_c.get(sev, DGREY)
        rect(s, Inches(0.4), y, Inches(1.1), Inches(0.38), clr)
        txt(s, code, Inches(0.4), y, Inches(1.1), Inches(0.38),
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(s, Inches(1.6), y, Inches(1.5), Inches(0.38), LGREY)
        txt(s, sev, Inches(1.6), y, Inches(1.5), Inches(0.38),
            size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(3.25), y, Inches(9.8), Inches(0.38), size=13, color=DGREY)
        y += Inches(0.46)

def render_mvno(prs, lang="EN"):
    t = "MVNO + opensecstack — Governed Network Operations" if lang=="EN" else "MVNO + opensecstack — Operacione Rrjeti të Qeverisura"
    s = slide_content(prs, t)
    events_en = [
        ("SIM Activation",  "Kerkese submitted → MARSHAL Gate 1 (SoD: activator ≠ approver)\n→ EXECUTE → WORM entry → eSIM provisioned"),
        ("Data Policy",     "Policy change → Gate 2 (scope: contract type) → Gate 4 (evidence: regulatory ref)\n→ EXECUTE → WORM logged → policy applied"),
        ("Roaming Config",  "Roaming enable → Gate 1 (authority) + Gate 3 (change ID required)\n→ EXECUTE → WORM → CDR hash stored"),
        ("VoIP Trunk",      "Trunk config → all 5 gates evaluated → EXECUTE + WORM\n→ Real-time VIGIL health check on trunk status"),
        ("Emergency Block", "SoD violation detected → HARD_STOP immediately\n→ WORM logged unconditionally → AUG-008 raised → ops alerted"),
    ]
    events_al = [
        ("Aktivizim SIM",   "Kerkese paraqitur → Porta 1 MARSHAL (NDS: aktivizues ≠ miratues)\n→ EXECUTE → hyrje WORM → eSIM provizionuar"),
        ("Politika e të dhënave", "Ndryshim politike → Porta 2 (fushëveprim: lloji i kontratës) → Porta 4 (prove)\n→ EXECUTE → WORM → politika zbatuar"),
        ("Konfig. Roaming",  "Aktivizim roaming → Porta 1 (autoritet) + Porta 3 (ID ndryshimi kërkohet)\n→ EXECUTE → WORM → hash CDR ruhet"),
        ("Trungu VoIP",      "Konfig. trungu → të 5 portat vlerësohen → EXECUTE + WORM\n→ Kontroll shëndetit VIGIL kohë reale mbi gjendjen e trungut"),
        ("Bllok Emergjence", "Shkelje NDS zbuluar → HARD_STOP menjëherë\n→ WORM i regjistruar pa kushte → AUG-008 ngritur → ops njoftuar"),
    ]
    events = events_en if lang=="EN" else events_al
    y = Inches(1.4)
    clrs = [BLUE, GREEN, TEAL, PURPLE, RED]
    for (ev, desc), clr in zip(events, clrs):
        rect(s, Inches(0.4), y, Inches(2.5), Inches(0.85), clr)
        txt(s, ev, Inches(0.45), y + Inches(0.08), Inches(2.4), Inches(0.7),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(s, Inches(3.05), y, Inches(10.0), Inches(0.85), LGREY)
        txt(s, desc, Inches(3.15), y + Inches(0.06), Inches(9.8), Inches(0.78),
            size=13, color=NAVY)
        y += Inches(1.02)

def render_nis2(prs, lang="EN"):
    t = "NIS2 Article 21(2) — Complete Compliance Matrix" if lang=="EN" else "NIS2 Neni 21(2) — Matrica e Plotë e Pajtueshmërisë"
    s = slide_content(prs, t)
    matrix_en = [
        ("(a) Risk analysis",       "L4",  "CITADEL governance + NIS2Compass art21_a"),
        ("(b) Incident handling",   "L4+5","WORM log + HARD_STOP playbook"),
        ("(c) Business continuity", "L5",  "Active-active + offsite DR + triple backup"),
        ("(d) Supply chain",        "L3",  "Dependency pinning + Go/Cargo checksums"),
        ("(e) Network security",    "L2+3","Rate limiting · TLS · Kubernetes isolation"),
        ("(f) Effectiveness",       "L4",  "APIGuard self-scan + VIGIL_DEEP daily audit"),
        ("(g) Cyber hygiene",       "L1",  "SoD enforcement + role model + training"),
        ("(h) Cryptography",        "L5",  "TripleHash + Ed25519 anchors + HMAC-SHA256"),
        ("(i) HR security",         "L1",  "Role groups + access revocation workflow"),
        ("(j) Access control",      "L1+4","JWT/API keys + MARSHAL gate evaluation"),
    ]
    matrix_al = [
        ("(a) Analiza rreziku",     "L4",  "Qeverisja CITADEL + NIS2Compass art21_a"),
        ("(b) Trajtimi incidentit", "L4+5","Regjistri WORM + libri i lojës HARD_STOP"),
        ("(c) Vazhdimësia biznesit","L5",  "Aktiv-aktiv + DR jashtë sitit + kopje triple"),
        ("(d) Zinxhiri furnizimit", "L3",  "Fiksim varësie + kontrolle Go/Cargo"),
        ("(e) Siguria rrjetit",     "L2+3","Kufizim rate · TLS · izolim Kubernetes"),
        ("(f) Efektiviteti",        "L4",  "Vetë-skanim APIGuard + auditim ditor VIGIL_DEEP"),
        ("(g) Higjiena kibernetike","L1",  "Zbatim NDS + model roli + trajnim"),
        ("(h) Kriptografia",        "L5",  "TripleHash + ankora Ed25519 + HMAC-SHA256"),
        ("(i) Siguria HR",          "L1",  "Grupet e roleve + rrjedhë revokimi aksesi"),
        ("(j) Kontrolli aksesit",   "L1+4","JWT/çelësat API + vlerësimi porta MARSHAL"),
    ]
    matrix = matrix_en if lang=="EN" else matrix_al
    hdr_en = ["NIS2 Measure", "Layer", "opensecstack Component"]
    hdr_al = ["Masa NIS2",    "Shtresa","Komponenti opensecstack"]
    hdr = hdr_en if lang=="EN" else hdr_al
    widths = [Inches(3.5), Inches(0.9), Inches(8.3)]
    xs = [Inches(0.3), Inches(3.9), Inches(4.9)]
    # header
    for h2, w2, x in zip(hdr, widths, xs):
        rect(s, x, Inches(1.35), w2, Inches(0.45), NAVY)
        txt(s, h2, x + Inches(0.05), Inches(1.38), w2, Inches(0.4),
            size=13, bold=True, color=WHITE)
    y = Inches(1.82)
    for i, (measure, layer, comp) in enumerate(matrix):
        bg = LGREY if i % 2 == 0 else WHITE
        for val, w2, x in zip([measure, layer, comp], widths, xs):
            rect(s, x, y, w2, Inches(0.47), bg)
            clr = RED if layer in ("L4+5","L1+4") else (GREEN if layer == "L5" else NAVY)
            txt(s, val, x + Inches(0.05), y + Inches(0.05), w2 - Inches(0.1), Inches(0.4),
                size=12, color=clr if val == layer else NAVY)
        y += Inches(0.47)

def render_cta(prs, lang="EN"):
    s = blank(prs)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, H - Inches(0.10), W, Inches(0.10), TEAL)
    rect(s, 0, 0, Inches(0.08), H, TEAL)
    t = "opensecstack v0.1.0 — Join Us" if lang=="EN" else "opensecstack v0.1.0 — Bashkohuni me Ne"
    sub = "Apache 2.0 · Open Source · Community-Driven" if lang=="EN" else "Apache 2.0 · Burim i Hapur · I Drejtuar nga Komuniteti"
    txt(s, t,   Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.3),
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.65),
        size=20, color=TEAL, align=PP_ALIGN.CENTER)

    road_en = [
        ("Phase 4 — 2026", "APIGuard v0.6.0 · IRFlow · ThreatFlow · CI/CD orchestration"),
        ("Phase 5 — 2027", "APIGuard v1.0.0 · OpenScrub · SecureLab · OpenCSIRT · OS Alpha"),
        ("v1.0 Production","Security audit · HA deployment · Third-party review"),
    ]
    road_al = [
        ("Faza 4 — 2026",   "APIGuard v0.6.0 · IRFlow · ThreatFlow · Orkestrim CI/CD"),
        ("Faza 5 — 2027",   "APIGuard v1.0.0 · OpenScrub · SecureLab · OpenCSIRT · OS Alfa"),
        ("v1.0 Prodhim",    "Auditim sigurie · Vendosje HA · Rishikim i palëve të treta"),
    ]
    road = road_en if lang=="EN" else road_al
    y = Inches(3.7)
    for phase, desc in road:
        rect(s, Inches(0.6), y, Inches(12.1), Inches(0.7), RGBColor(0x14, 0x2A, 0x44))
        txt(s, phase, Inches(0.75), y + Inches(0.1), Inches(2.8), Inches(0.55),
            size=15, bold=True, color=TEAL)
        txt(s, desc,  Inches(3.6),  y + Inches(0.1), Inches(9.0), Inches(0.55),
            size=15, color=WHITE)
        y += Inches(0.82)

    pillars_en = [("Immutable","WORM chain"), ("Governed","MARSHAL 5-gate"),
                  ("Auditable","VIGIL_DEEP"), ("Open","Apache 2.0")]
    pillars_al = [("I pandryshueshëm","Zinxhir WORM"), ("I qeverisur","MARSHAL 5 porta"),
                  ("I auditushëm","VIGIL_DEEP"),   ("I hapur","Apache 2.0")]
    pillars = pillars_en if lang=="EN" else pillars_al
    for i, (lbl, desc) in enumerate(pillars):
        x = Inches(0.6 + i * 3.1)
        rect(s, x, Inches(6.2), Inches(2.9), Inches(0.95), TEAL)
        txt(s, lbl,  x, Inches(6.25), Inches(2.9), Inches(0.45),
            size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(s, desc, x, Inches(6.7), Inches(2.9), Inches(0.4),
            size=13, color=NAVY, align=PP_ALIGN.CENTER)

    txt(s, "github.com/opensecstack/opensecstack  ·  Questions?",
        Inches(0), H - Inches(0.52), W, Inches(0.42),
        size=14, italic=True, color=LGREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# RENDERER DISPATCH
# ══════════════════════════════════════════════════════════════════════════════

CUSTOM = {
    "custom_ecosystem":      lambda p: render_ecosystem(p, "EN"),
    "custom_ecosystem_al":   lambda p: render_ecosystem(p, "AL"),
    "custom_citadel":        lambda p: render_citadel(p, "EN"),
    "custom_citadel_al":     lambda p: render_citadel(p, "AL"),
    "custom_layers":         lambda p: render_layers(p, "EN"),
    "custom_layers_al":      lambda p: render_layers(p, "AL"),
    "custom_microkernel":    lambda p: render_microkernel(p, "EN"),
    "custom_microkernel_al": lambda p: render_microkernel(p, "AL"),
    "custom_desktop_os":     lambda p: render_desktop_os(p, "EN"),
    "custom_desktop_os_al":  lambda p: render_desktop_os(p, "AL"),
    "custom_mobile_os":      lambda p: render_mobile_os(p, "EN"),
    "custom_mobile_os_al":   lambda p: render_mobile_os(p, "AL"),
    "custom_sha256":         lambda p: render_hash(p, "sha256", "EN"),
    "custom_sha256_al":      lambda p: render_hash(p, "sha256", "AL"),
    "custom_sha512":         lambda p: render_hash(p, "sha512", "EN"),
    "custom_sha512_al":      lambda p: render_hash(p, "sha512", "AL"),
    "custom_blake3":         lambda p: render_hash(p, "blake3", "EN"),
    "custom_blake3_al":      lambda p: render_hash(p, "blake3", "AL"),
    "custom_triplehash":     lambda p: render_triplehash(p, "EN"),
    "custom_triplehash_al":  lambda p: render_triplehash(p, "AL"),
    "custom_tds":            lambda p: render_tds(p, "EN"),
    "custom_tds_al":         lambda p: render_tds(p, "AL"),
    "custom_marshal":        lambda p: render_marshal(p, "EN"),
    "custom_marshal_al":     lambda p: render_marshal(p, "AL"),
    "custom_augur":          lambda p: render_augur(p, "EN"),
    "custom_augur_al":       lambda p: render_augur(p, "AL"),
    "custom_mvno":           lambda p: render_mvno(p, "EN"),
    "custom_mvno_al":        lambda p: render_mvno(p, "AL"),
    "custom_nis2":           lambda p: render_nis2(p, "EN"),
    "custom_nis2_al":        lambda p: render_nis2(p, "AL"),
    "custom_cta":            lambda p: render_cta(p, "EN"),
    "custom_cta_al":         lambda p: render_cta(p, "AL"),
}


def build_deck(lang_key):
    cfg = DECKS[lang_key]
    prs = new_prs()
    # Title slide
    slide_title(prs, cfg["title"], cfg["subtitle"], cfg["footer"])
    # Content slides
    for entry in cfg["slides"]:
        stype = entry[0]
        if stype == "section":
            slide_section(prs, entry[1], entry[2], entry[3] if len(entry) > 3 else "")
        elif stype == "bullets":
            bullets(prs, entry[1], entry[2])
        elif stype.startswith("custom_"):
            CUSTOM[stype](prs)
    return prs


# ── Build & save ──────────────────────────────────────────────────────────────
outdir = os.path.dirname(__file__)
for lang in ("EN", "AL"):
    prs = build_deck(lang)
    path = os.path.join(outdir, DECKS[lang]["filename"])
    prs.save(path)
    print(f"Saved [{lang}]: {path}  ({len(prs.slides)} slides)")
