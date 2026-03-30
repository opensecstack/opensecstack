"""
opensecstack conference deck generator.
Produces opensecstack_conference.pptx in the same directory.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand palette ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy — background / title
TEAL    = RGBColor(0x00, 0xC2, 0xB2)   # bright teal — accent / headings
AMBER   = RGBColor(0xFF, 0xB7, 0x00)   # amber — warnings / highlight
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGREY   = RGBColor(0xE8, 0xEC, 0xF0)   # light grey — body text bg
DGREY   = RGBColor(0x44, 0x55, 0x66)   # dark grey — body text
GREEN   = RGBColor(0x00, 0xCC, 0x88)   # green — PASS / ok
RED     = RGBColor(0xFF, 0x3B, 0x3B)   # red — FAIL / critical

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.line.color.rgb = fill
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, x, y, w, h,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def navy_slide(title_text, subtitle_text=None):
    """Dark navy full-bleed slide — used for section breaks and title."""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, NAVY)
    # teal accent bar bottom
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), TEAL)
    # title
    add_text(slide, title_text,
             Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.8),
             size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle_text:
        add_text(slide, subtitle_text,
                 Inches(1.2), Inches(4.1), Inches(11.0), Inches(1.0),
                 size=22, color=TEAL, align=PP_ALIGN.CENTER)
    return slide


def content_slide(title_text):
    """White content slide with navy header bar."""
    slide = prs.slides.add_slide(BLANK)
    # white bg
    add_rect(slide, 0, 0, W, H, WHITE)
    # navy header
    add_rect(slide, 0, 0, W, Inches(1.15), NAVY)
    # teal left accent strip
    add_rect(slide, 0, Inches(1.15), Inches(0.08), H - Inches(1.15), TEAL)
    # title text
    add_text(slide, title_text,
             Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.85),
             size=28, bold=True, color=WHITE)
    return slide


def bullet_slide(title, bullets):
    """Standard content slide with bullet list."""
    slide = content_slide(title)
    y = Inches(1.4)
    for b in bullets:
        indent = b.startswith("  ")
        txt  = b.lstrip()
        sz   = 18 if indent else 21
        clr  = DGREY if indent else NAVY
        pfx  = "    ◦  " if indent else "  ▸  "
        add_text(slide, pfx + txt,
                 Inches(0.5), y, Inches(12.3), Inches(0.5),
                 size=sz, color=clr, bold=not indent)
        y += Inches(0.52 if not indent else 0.44)
    return slide


def add_box(slide, label, value, x, y, w=Inches(2.8), h=Inches(1.5),
            bg=NAVY, label_clr=TEAL, val_clr=WHITE):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, label, x + Inches(0.1), y + Inches(0.1),
             w - Inches(0.2), Inches(0.5), size=13, color=label_clr, bold=True)
    add_text(slide, value, x + Inches(0.1), y + Inches(0.5),
             w - Inches(0.2), Inches(0.85), size=18, color=val_clr, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = navy_slide(
    "opensecstack",
    "Immutable Governance · Blockchain-Anchored Audit · Ecosystem Security"
)
add_text(s, "Conference Presentation  ·  2026",
         Inches(0), H - Inches(0.55), W, Inches(0.45),
         size=14, color=LGREY, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
bullet_slide("The Problem — Why Governance Fails", [
    "AI systems and automated pipelines execute thousands of actions per day",
    "  No verifiable record of who authorised what, when, and why",
    "Separation of duties (SoD) is bypassed silently in complex systems",
    "  Same operator triggers and approves the same change",
    "Audit logs are mutable — stored in databases that admins can edit",
    "  Post-breach forensics cannot be trusted",
    "Regulatory mandates (NIS2, ISO 27001) require evidence of control",
    "  'We had a process' is not a compliant answer",
])

# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Ecosystem Overview
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("The opensecstack Ecosystem")
# Draw 4 platform boxes
boxes = [
    ("APIGuard",    "API Security\nOWASP API Top 10\nCVSS 3.1 · SARIF",    Inches(0.5),  Inches(1.5)),
    ("NIS2 Compass","Regulatory Compliance\nNIS2 Article 21(2)\n10 Measures A–J", Inches(3.7),  Inches(1.5)),
    ("CITADEL",     "Governance Engine\nMARSHAL · WORM · VIGIL\nAUGUR · Anchor", Inches(6.9),  Inches(1.5)),
    ("SDK",         "Multi-language\nGo · Python · Rust\nTDS Scanner",      Inches(10.1), Inches(1.5)),
]
for lbl, val, x, y in boxes:
    add_box(s, lbl, val, x, y, w=Inches(2.9), h=Inches(1.9))

# Arrow connectors (simple rectangles)
for ax in [Inches(3.4), Inches(6.6), Inches(9.8)]:
    add_rect(s, ax, Inches(2.3), Inches(0.3), Inches(0.08), TEAL)

# bottom connector label
add_text(s, "All platforms emit signed WORM events to CITADEL via HMAC-SHA256 authenticated connectors",
         Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.6),
         size=16, color=DGREY, italic=True, align=PP_ALIGN.CENTER)

# Stats row
stats = [
    ("5",    "OWASP Gates\n(MARSHAL)"),
    ("10",   "NIS2 Measures\n(Compass)"),
    ("3",    "SDK Languages"),
    ("SHA-256", "Chain Hash\nAlgorithm"),
    ("Ed25519", "Anchor\nSignature"),
]
for i, (val, lbl) in enumerate(stats):
    bx = Inches(0.5 + i * 2.55)
    add_rect(s, bx, Inches(4.5), Inches(2.35), Inches(1.5), TEAL)
    add_text(s, val, bx, Inches(4.55), Inches(2.35), Inches(0.7),
             size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, lbl, bx, Inches(5.25), Inches(2.35), Inches(0.65),
             size=12, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 — APIGuard
# ══════════════════════════════════════════════════════════════════════════════
bullet_slide("APIGuard — API Security Testing", [
    "Scans OpenAPI 3.x, Swagger 2.0, GraphQL, and Postman collections",
    "Full OWASP API Security Top 10 (A1–A10) coverage",
    "  A1 BOLA · A2 Broken Auth · A3 Mass Assignment · A4 Rate Limiting",
    "  A5 Function Auth · A6 Sensitive Flows · A7 SSRF · A8–A10 + more",
    "CVSS 3.1 scoring on every finding",
    "Output: SARIF, HTML, JSON, plain-text reports",
    "React dashboard — scan history, finding management, regression detection",
    "Custom rule engine — load organisation-specific checks as WASM modules",
])

# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — NIS2 Compass
# ══════════════════════════════════════════════════════════════════════════════
bullet_slide("NIS2 Compass — Regulatory Compliance", [
    "Maps all 10 NIS2 Article 21(2) technical measures (A through J)",
    "  A: Risk analysis  ·  B: Incident handling  ·  C: Business continuity",
    "  D: Supply chain  ·  E: Acquisition security  ·  F: Vulnerability mgmt",
    "  G: Policies & procedures  ·  H: Cryptography  ·  I: HR security  ·  J: Access control",
    "Evidence artefacts linked to each control — cryptographic hashes stored",
    "Assessment workflows with verifier sign-off",
    "CITADEL webhook integration — compliance state changes → VIGIL advisories",
    "PDF + JSON reports for regulator submission",
])

# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 — CITADEL Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("CITADEL — The Governance Engine")
components = [
    ("MARSHAL",  "5-Gate\nDecision Engine",   Inches(0.4),  Inches(1.6), NAVY),
    ("WORM Log", "Append-Only\nAudit Chain",  Inches(3.15), Inches(1.6), NAVY),
    ("VIGIL",    "Real-Time\nHealth Monitor", Inches(5.9),  Inches(1.6), NAVY),
    ("AUGUR",    "Predictive\nAdvisory Rules",Inches(8.65), Inches(1.6), NAVY),
    ("Anchor",   "Ed25519\nChain Rotation",   Inches(11.4), Inches(1.6), NAVY),
]
for lbl, val, x, y, bg in components:
    add_box(s, lbl, val, x, y, w=Inches(2.55), h=Inches(1.6), bg=bg)

add_text(s, "⟵────────────────── All components write to the WORM log ──────────────────⟶",
         Inches(0.3), Inches(3.45), Inches(12.7), Inches(0.5),
         size=14, color=TEAL, align=PP_ALIGN.CENTER, italic=True)

bullets2 = [
    "Connector auth: HMAC-SHA256 signed headers (X-CITADEL-KEY · X-CITADEL-TS · X-CITADEL-SIG)",
    "Admin auth: HS256 JWT with role-based access (group_sig_admin · group_sig_auditor)",
    "Two-bucket sliding window replay protection on every connector request",
    "No mutable state — all decisions are final and chain-linked",
]
y = Inches(4.1)
for b in bullets2:
    add_text(s, "  ▸  " + b, Inches(0.5), y, Inches(12.3), Inches(0.5),
             size=17, color=DGREY)
    y += Inches(0.52)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — MARSHAL 5-Gate Engine
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("MARSHAL — The 5-Gate Decision Engine")
gates = [
    ("Gate 1", "AUTHORITY",     "SoD check\nOperator ≠ Verifier\nRole validation",       GREEN),
    ("Gate 2", "SCOPE",         "Project whitelist\nFreeze detection\nAction type allow",  TEAL),
    ("Gate 3", "DETERMINISM",   "Change ID required\nTrace to ticket\nNo orphan actions", AMBER),
    ("Gate 4", "EVIDENCE",      "Artefact hashes\nChange ID linkage\nAudit artifacts",    RGBColor(0x88,0x44,0xFF)),
    ("Gate 5", "SCHEMA",        "Kerkese 2.0\nJSON Schema\nField validation",             RGBColor(0xFF,0x66,0x00)),
]
for i, (g, name, desc, clr) in enumerate(gates):
    x = Inches(0.25 + i * 2.6)
    add_rect(s, x, Inches(1.4), Inches(2.45), Inches(3.8), clr)
    add_text(s, g,    x, Inches(1.45), Inches(2.45), Inches(0.5),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, name, x, Inches(1.9),  Inches(2.45), Inches(0.6),
             size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x, Inches(2.55), Inches(2.45), Inches(2.4),
             size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Outcome: EXECUTE   ·   REFUSE   ·   HARD_STOP  — written to WORM log on every decision",
         Inches(0.3), Inches(5.45), Inches(12.7), Inches(0.5),
         size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "HARD_STOP is never skipped, even on internal error. The immutable record is the source of truth.",
         Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.5),
         size=14, italic=True, color=DGREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 — WORM Log
# ══════════════════════════════════════════════════════════════════════════════
bullet_slide("WORM Log — Write Once, Read Many", [
    "Every governance event is appended — never updated, never deleted",
    "  Scan decisions · Finding evidence · Freeze actions · Key rotations",
    "Each entry carries: entry_id, source, event_type, project_id, payload, chain_hash, prev_hash",
    "Chain hash links every entry to its predecessor — break one, break all",
    "Anchor records (Ed25519 signatures) are themselves WORM entries",
    "Full chain verification via POST /api/v1/worm/verify",
    "  Returns: entries_checked, chain_integrity PASS/FAIL, first_break_entry",
    "Genesis hash: 64 zero bytes — mathematically provable chain origin",
])

# ══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Blockchain-Inspired Design
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("Blockchain-Inspired Chain Design")

# Draw chain diagram
chain_items = [
    ("GENESIS\n0000…0000", DGREY),
    ("Entry #1\nMARSHAL\nEXECUTE", NAVY),
    ("Entry #2\nScan\nComplete", NAVY),
    ("Entry #3\nAnchor\nRotation", RGBColor(0x44,0x00,0xCC)),
    ("Entry N\n…", NAVY),
]
for i, (lbl, clr) in enumerate(chain_items):
    bx = Inches(0.35 + i * 2.5)
    add_rect(s, bx, Inches(2.0), Inches(2.1), Inches(1.8), clr)
    add_text(s, lbl, bx, Inches(2.1), Inches(2.1), Inches(1.7),
             size=14, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
    if i < len(chain_items) - 1:
        add_rect(s, bx + Inches(2.1), Inches(2.8), Inches(0.4), Inches(0.1), TEAL)

add_text(s, "prev_hash →", Inches(0.35), Inches(3.95), Inches(2.1), Inches(0.4),
         size=12, color=TEAL, align=PP_ALIGN.CENTER)
for i in range(1, 5):
    bx = Inches(0.35 + i * 2.5)
    add_text(s, "chain_hash\nprev=hash(i-1)", bx, Inches(3.95), Inches(2.1), Inches(0.6),
             size=11, color=TEAL, align=PP_ALIGN.CENTER)

props = [
    ("Immutability",    "Changing any entry invalidates every subsequent hash"),
    ("Auditability",    "Any verifier can replay the chain from genesis independently"),
    ("Tamper Evidence", "A single bit flip is detectable without a trusted third party"),
    ("Anchoring",       "Ed25519 signatures periodically seal the chain tip"),
]
y = Inches(4.9)
for prop, desc in props:
    add_text(s, f"▸  {prop}: ", Inches(0.5), y, Inches(2.8), Inches(0.45),
             size=16, bold=True, color=NAVY)
    add_text(s, desc, Inches(3.3), y, Inches(9.5), Inches(0.45),
             size=16, color=DGREY)
    y += Inches(0.48)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Triple Hashing
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("Triple Hashing — The Chain Integrity Model")

add_text(s, "Three independent SHA-256 hash operations protect every decision:",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         size=18, color=DGREY)

hashes = [
    ("Hash 1", "PAYLOAD HASH",
     "SHA-256(canonical_json(payload))",
     "Ensures payload cannot be altered post-write.\nCanonical JSON: sorted keys, no whitespace.",
     NAVY),
    ("Hash 2", "CHAIN HASH",
     "SHA-256(prev_chain_hash || payload_hash)",
     "Links this entry to every preceding entry.\nBreaking the chain requires rewriting all successors.",
     RGBColor(0x00,0x55,0xAA)),
    ("Hash 3", "BODY HASH (Connector Auth)",
     "SHA-256(request_body)",
     "Prevents in-transit tampering of the API request\nbefore it reaches the WORM appender.",
     RGBColor(0x44,0x00,0xCC)),
]
for i, (num, name, formula, desc, clr) in enumerate(hashes):
    y = Inches(2.1 + i * 1.65)
    add_rect(s, Inches(0.4), y, Inches(12.5), Inches(1.5), RGBColor(0xF0,0xF4,0xF8))
    add_rect(s, Inches(0.4), y, Inches(0.1), Inches(1.5), clr)
    add_text(s, num,     Inches(0.6), y + Inches(0.08), Inches(1.0), Inches(0.4),
             size=13, bold=True, color=clr)
    add_text(s, name,    Inches(1.6), y + Inches(0.08), Inches(5.0), Inches(0.4),
             size=16, bold=True, color=NAVY)
    add_text(s, formula, Inches(6.7), y + Inches(0.05), Inches(6.1), Inches(0.45),
             size=13, color=RGBColor(0x22,0x44,0x88), italic=True)
    add_text(s, desc,    Inches(1.6), y + Inches(0.55), Inches(11.1), Inches(0.8),
             size=14, color=DGREY)

add_text(s, "Together these three hashes make it cryptographically impossible to forge, replay, or silently modify any governance record.",
         Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.45),
         size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 11 — Chain Anchoring
# ══════════════════════════════════════════════════════════════════════════════
bullet_slide("Chain Anchoring — Ed25519 Signatures", [
    "Problem: SHA-256 chains prove internal consistency but not calendar time",
    "  An attacker with enough compute could reconstruct a plausible chain",
    "Solution: periodic Ed25519 anchor signatures over the current chain tip",
    "  Anchor message: 'anchor:' + current_chain_hash",
    "  Anchor record written as a WORM entry — anchors are chained too",
    "Key management: key ID + public key stored separately from private key",
    "  Rotation scheduler runs every N hours (configurable, default 720h)",
    "  Old anchors remain verifiable — public keys are append-only",
    "Verification: anyone with the public key can verify without server access",
    "  Ed25519 verify(pubkey, 'anchor:' + chain_hash, signature)",
])

# ══════════════════════════════════════════════════════════════════════════════
# Slide 12 — VIGIL
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("VIGIL — Real-Time Health Monitoring")

states = [
    ("GREEN",  "Normal operation\nNo active HARD_STOPs\nAll mirrors fresh",        GREEN),
    ("AMBER",  "Elevated risk\nHIGH advisory active\nMirror stale > threshold",    AMBER),
    ("RED",    "Immediate action required\nCRITICAL advisory\nor active HARD_STOP", RED),
]
for i, (st, desc, clr) in enumerate(states):
    x = Inches(0.5 + i * 4.2)
    add_rect(s, x, Inches(1.4), Inches(3.9), Inches(2.2), clr)
    add_text(s, st,  x, Inches(1.5),  Inches(3.9), Inches(0.65),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x, Inches(2.15), Inches(3.9), Inches(1.3),
             size=15, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "VIGIL_REALTIME — polled every N seconds, cached, zero DB load on status requests",
         Inches(0.4), Inches(3.8), Inches(12.5), Inches(0.45),
         size=15, italic=True, color=DGREY, align=PP_ALIGN.CENTER)

deep = [
    ("VIGIL_DEEP — Daily Full-Chain Audit", ""),
    ("Chain integrity",   "Walk every WORM entry, verify prev_hash → chain_hash"),
    ("SoD violations",    "Detect entries where operator_user_id == verifier_user_id"),
    ("Orphaned incidents","HARD_STOP entries with no subsequent close_incident EXECUTE"),
    ("Evidence gaps",     "EXECUTE decisions with no linked artefacts or change ID"),
    ("Signed report",     "SHA-256(report_id|period|summary) written as a WORM entry"),
]
y = Inches(4.4)
for lbl, val in deep:
    if not val:
        add_text(s, lbl, Inches(0.5), y, Inches(12.3), Inches(0.45),
                 size=17, bold=True, color=NAVY)
    else:
        add_text(s, f"  ▸  {lbl}:", Inches(0.5), y, Inches(4.0), Inches(0.4),
                 size=14, bold=True, color=TEAL)
        add_text(s, val, Inches(4.5), y, Inches(8.3), Inches(0.4),
                 size=14, color=DGREY)
    y += Inches(0.43)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 13 — AUGUR
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("AUGUR — Predictive Advisory Engine")
add_text(s, "9 advisory rules evaluated continuously against live system state:",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.45),
         size=17, color=DGREY)

advisories = [
    ("AUG-001", "MEDIUM",   "Mandate expiry approaching — renew before MARSHAL Gate 1 refuses"),
    ("AUG-002", "HIGH",     "Approval threshold exceeded — additional sign-off required"),
    ("AUG-003", "MEDIUM",   "Vendor contract expiry — supply chain risk elevation"),
    ("AUG-004", "HIGH",     "Freeze period active — MARSHAL Gate 2 will REFUSE non-emergency"),
    ("AUG-005", "MEDIUM",   "Mirror data inconsistency — cross-check ERP before proceeding"),
    ("AUG-006", "MEDIUM",   "Unusual action frequency — possible anomalous automation"),
    ("AUG-007", "HIGH",     "Mirror stale beyond threshold — advisory data may be outdated"),
    ("AUG-008", "CRITICAL", "Unresolved HARD_STOP — all non-emergency actions at elevated risk"),
    ("AUG-009", "HIGH",     "Chain anchor overdue — rotation required, timestamp trust weakens"),
]
sev_color = {"CRITICAL": RED, "HIGH": AMBER, "MEDIUM": TEAL, "LOW": GREEN}
y = Inches(1.9)
for code, sev, desc in advisories:
    clr = sev_color.get(sev, DGREY)
    add_rect(s, Inches(0.4), y, Inches(1.0), Inches(0.36), clr)
    add_text(s, code, Inches(0.4), y, Inches(1.0), Inches(0.36),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.5), y, Inches(1.4), Inches(0.36), RGBColor(0xEE,0xEE,0xEE))
    add_text(s, sev, Inches(1.5), y, Inches(1.4), Inches(0.36),
             size=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(3.0), y, Inches(10.0), Inches(0.36),
             size=13, color=DGREY)
    y += Inches(0.44)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 14 — Full Ecosystem Flow
# ══════════════════════════════════════════════════════════════════════════════
s = content_slide("End-to-End Governance Flow")

steps = [
    ("1", "Developer triggers\nAPI scan via\nAPIGuard CLI",          Inches(0.3)),
    ("2", "APIGuard sends\nKerkese to CITADEL\nMARSHAL evaluate",    Inches(2.85)),
    ("3", "MARSHAL runs\n5 gates\nIssues EXECUTE",                   Inches(5.4)),
    ("4", "Scan runs\nFindings written to\nWORM as evidence",        Inches(7.95)),
    ("5", "VIGIL reflects\nnew state\nAUGUR evaluates",              Inches(10.5)),
]
colors = [NAVY, RGBColor(0x00,0x55,0xAA), RGBColor(0x00,0x88,0x66), NAVY, RGBColor(0x44,0x00,0xCC)]
for (num, lbl, x), clr in zip(steps, colors):
    add_rect(s, x, Inches(1.5), Inches(2.4), Inches(2.0), clr)
    add_text(s, num, x + Inches(0.05), Inches(1.55), Inches(0.5), Inches(0.5),
             size=22, bold=True, color=TEAL)
    add_text(s, lbl, x, Inches(2.0), Inches(2.4), Inches(1.4),
             size=13, color=WHITE, align=PP_ALIGN.CENTER)
    if x < Inches(10.5):
        add_rect(s, x + Inches(2.4), Inches(2.4), Inches(0.45), Inches(0.1), TEAL)

add_text(s, "Every step is recorded in the immutable WORM chain. No step can be skipped, forged, or silently failed.",
         Inches(0.3), Inches(3.75), Inches(12.7), Inches(0.5),
         size=15, italic=True, color=DGREY, align=PP_ALIGN.CENTER)

noncompliant = [
    ("Without opensecstack",  "No governance record · SoD bypass undetected · Audit logs mutable · NIS2 evidence absent", RED),
    ("With opensecstack",     "Every decision chain-hashed · SoD enforced at Gate 1 · WORM log forensically sound · NIS2 mapped",    GREEN),
]
y = Inches(4.45)
for lbl, desc, clr in noncompliant:
    add_rect(s, Inches(0.3), y, Inches(0.12), Inches(0.7), clr)
    add_text(s, lbl + ": ", Inches(0.55), y + Inches(0.05), Inches(3.1), Inches(0.6),
             size=15, bold=True, color=NAVY)
    add_text(s, desc, Inches(3.7), y + Inches(0.1), Inches(9.2), Inches(0.55),
             size=14, color=DGREY)
    y += Inches(0.85)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 15 — Summary & Call to Action
# ══════════════════════════════════════════════════════════════════════════════
s = navy_slide(
    "opensecstack v0.1.0",
    "Open-source · Apache 2.0 · Production-ready governance"
)
pillars = [
    ("Immutable",   "WORM chain — SHA-256\nlinked, append-only,\nforensically sound"),
    ("Governed",    "MARSHAL 5-gate engine —\nevery action authorised\nbefore execution"),
    ("Auditable",   "VIGIL_DEEP daily audit —\nchain integrity + SoD +\nevidence gaps"),
    ("Anchored",    "Ed25519 signatures —\nperiodic seal of the\nentire chain history"),
]
for i, (title, desc) in enumerate(pillars):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(4.9), Inches(3.0), Inches(1.9), TEAL)
    add_text(s, title, x, Inches(4.95), Inches(3.0), Inches(0.55),
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, desc,  x, Inches(5.5),  Inches(3.0), Inches(1.2),
             size=13, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s, "github.com/opensecstack/opensecstack",
         Inches(0), H - Inches(0.5), W, Inches(0.4),
         size=14, color=LGREY, align=PP_ALIGN.CENTER, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(__file__), "opensecstack_conference.pptx")
prs.save(out)
print(f"Saved: {out}")
